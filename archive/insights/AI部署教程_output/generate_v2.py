#!/usr/bin/env python3
"""
AI 部署教程 — PDF + PPT + Excel 一站式生成器
覆盖：Claude Code / DeepSeek / Codex / Cherry Studio
"""
import os, sys
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent
IMG_DIR = OUTPUT_DIR / "images"
IMG_DIR.mkdir(exist_ok=True)

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

C = {
    "primary": "#4F46E5", "success": "#059669", "warning": "#D97706",
    "danger": "#DC2626", "info": "#0891B2", "dark": "#1E293B",
    "accent": "#7C3AED", "bg": "#F8FAFC", "ds": "#1A56DB",
}

# ============================================================
# 新增图表：DeepSeek vs Claude 双路线对比
# ============================================================
def fig_deepseek_vs_claude():
    fig, axes = plt.subplots(1, 2, figsize=(11, 5.5))
    fig.suptitle("🚀 两条部署路线对比：国内用户首选 DeepSeek", fontsize=14, fontweight="bold", color=C["dark"], y=0.98)
    fig.set_facecolor("white")

    # 左图：Claude 路线
    ax = axes[0]
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")
    ax.set_title("🌍 路线A：Claude Code + Anthropic", fontsize=12, fontweight="bold", color=C["primary"])
    claude_steps = [
        (5, 8.5, "注册 Anthropic", "需外网 + 海外手机号", C["warning"]),
        (5, 6.5, "海外支付 $5", "需 Visa/MasterCard", C["danger"]),
        (5, 4.5, "获取 API Key", "console.anthropic.com", C["primary"]),
        (5, 2.5, "安装 Claude Code", "npm install -g ...", C["accent"]),
    ]
    for x, y, title, desc, color in claude_steps:
        box = FancyBboxPatch((x-2.5, y-0.6), 5, 1.2, boxstyle="round,pad=0.1",
                             facecolor=color, edgecolor="white", linewidth=1.5, alpha=0.85)
        ax.add_patch(box)
        ax.text(x, y+0.15, title, ha="center", fontsize=10, fontweight="bold", color="white")
        ax.text(x, y-0.3, desc, ha="center", fontsize=7.5, color="white", alpha=0.9)
    # 箭头
    for i in range(3):
        ax.annotate("", xy=(5, claude_steps[i+1][1]+0.7), xytext=(5, claude_steps[i][1]-0.7),
                    arrowprops=dict(arrowstyle="->", color=C["dark"], lw=1.5))
    # 难度标签
    ax.text(5, 0.5, "⚠️ 门槛较高：需要外网 + 海外支付", ha="center", fontsize=9, color=C["danger"], fontweight="bold")

    # 右图：DeepSeek 路线
    ax = axes[1]
    ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")
    ax.set_title("🇨🇳 路线B：DeepSeek API + Cherry Studio（推荐）", fontsize=12, fontweight="bold", color=C["ds"])
    ds_steps = [
        (5, 8.5, "注册 DeepSeek", "platform.deepseek.com\n国内直接访问 ✅", C["success"]),
        (5, 6.5, "微信/支付宝充值", "最低 ¥10 起充 ✅", C["success"]),
        (5, 4.5, "获取 API Key", "控制台 → API Keys", C["ds"]),
        (5, 2.5, "安装 Cherry Studio", "cherry-ai.com 下载\n填入 Key 即可用", C["info"]),
    ]
    for x, y, title, desc, color in ds_steps:
        box = FancyBboxPatch((x-2.5, y-0.65), 5, 1.3, boxstyle="round,pad=0.1",
                             facecolor=color, edgecolor="white", linewidth=1.5, alpha=0.85)
        ax.add_patch(box)
        ax.text(x, y+0.15, title, ha="center", fontsize=10, fontweight="bold", color="white")
        ax.text(x, y-0.35, desc, ha="center", fontsize=7.5, color="white", alpha=0.9)
    for i in range(3):
        ax.annotate("", xy=(5, ds_steps[i+1][1]+0.75), xytext=(5, ds_steps[i][1]-0.75),
                    arrowprops=dict(arrowstyle="->", color=C["dark"], lw=1.5))
    ax.text(5, 0.5, "✅ 零门槛：国内网络 + 微信支付宝即可", ha="center", fontsize=9, color=C["success"], fontweight="bold")

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    path = IMG_DIR / "deepseek_vs_claude.png"
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white", edgecolor="none")
    plt.close(fig)
    print(f"  ✅ {path.name}")
    return path

# ============================================================
# 如果已有图表则跳过重新生成
# ============================================================
def ensure_charts():
    charts = {}
    needed = {
        'roadmap': IMG_DIR / "deployment_roadmap.png",
        'compare': IMG_DIR / "tool_comparison.png",
        'trouble': IMG_DIR / "troubleshooting.png",
        'env': IMG_DIR / "env_checklist.png",
        'steps': IMG_DIR / "install_steps.png",
        'ds_vs_claude': IMG_DIR / "deepseek_vs_claude.png",
    }
    # Always regenerate ds_vs_claude (new chart)
    charts['ds_vs_claude'] = fig_deepseek_vs_claude()
    # Check existing
    for key, path in needed.items():
        if key == 'ds_vs_claude':
            continue
        if path.exists():
            charts[key] = path
        else:
            print(f"  ⚠️ Missing chart: {path.name}, regenerating all...")
            # Regenerate all (fallback)
            import subprocess
            subprocess.run([sys.executable, "generate_tutorial.py"], cwd=str(OUTPUT_DIR))
            for k, p in needed.items():
                if p.exists():
                    charts[k] = p
            break
    return charts

# ============================================================
# PDF 生成（fpdf2，紧凑排版）
# ============================================================
def build_pdf(charts):
    from fpdf import FPDF

    class PDF(FPDF):
        def __init__(self):
            super().__init__("P", "mm", "A4")
            font_dir = "C:/Windows/Fonts"
            for fn in ["msyh.ttc", "msyh.ttf", "simhei.ttf"]:
                fp = os.path.join(font_dir, fn)
                if os.path.exists(fp):
                    self.add_font("CJK", "", fp)
                    self.add_font("CJK", "B", fp)
                    break
            self.set_auto_page_break(True, 15)

        def h1(self, txt):
            self.ln(3)
            self.set_fill_color(*hex_rgb(C["primary"]))
            self.rect(self.l_margin, self.get_y(), 2.5, 8, "F")
            self.set_font("CJK", "B", 17)
            self.set_text_color(*hex_rgb(C["primary"]))
            self.set_x(self.l_margin + 4)
            self.cell(0, 8, txt, new_x="LMARGIN", new_y="NEXT")
            self.ln(2)

        def h2(self, txt):
            self.ln(1)
            self.set_font("CJK", "B", 12)
            self.set_text_color(*hex_rgb(C["dark"]))
            self.cell(0, 6, txt, new_x="LMARGIN", new_y="NEXT")
            self.set_draw_color(226, 232, 240)
            self.line(self.l_margin, self.get_y()+0.5, self.w-self.r_margin, self.get_y()+0.5)
            self.ln(2)

        def h3(self, txt):
            self.set_font("CJK", "B", 10.5)
            self.set_text_color(51, 65, 85)
            self.cell(0, 5.5, txt, new_x="LMARGIN", new_y="NEXT")
            self.ln(0.5)

        def p(self, txt):
            self.set_font("CJK", "", 9.5)
            self.set_text_color(*hex_rgb(C["dark"]))
            self.multi_cell(0, 4.8, txt, align="L")
            self.ln(0.5)

        def tip(self, txt, border=(59,130,246)):
            self.set_fill_color(*hex_rgb(C["bg"]))
            self.set_draw_color(*border)
            self.set_font("CJK", "", 8.5)
            self.set_text_color(*hex_rgb(C["dark"]))
            self.set_x(self.l_margin + 1)
            self.multi_cell(self.w - self.l_margin - self.r_margin - 2, 4.2, txt, fill=True)
            self.ln(1)

        def code(self, txt):
            self.set_fill_color(30, 41, 59)
            self.set_text_color(226, 232, 240)
            self.set_font("CJK", "", 7.8)
            self.set_x(self.l_margin + 1)
            self.multi_cell(self.w - self.l_margin - self.r_margin - 2, 3.8, txt, fill=True)
            self.set_text_color(*hex_rgb(C["dark"]))
            self.ln(1)

        def img(self, path, w=170):
            if path and Path(path).exists():
                self.image(str(path), x=self.l_margin + 3, w=w)
                self.ln(1.5)

        def tbl(self, headers, rows, widths=None):
            if widths is None:
                widths = [(self.w - self.l_margin - self.r_margin) / len(headers)] * len(headers)
            self.set_fill_color(*hex_rgb(C["primary"]))
            self.set_text_color(255, 255, 255)
            self.set_font("CJK", "B", 8)
            for h, w in zip(headers, widths):
                self.cell(w, 5.5, h, border=0, fill=True, align="L")
            self.ln()
            for i, row in enumerate(rows):
                self.set_fill_color(*hex_rgb(C["bg"] if i%2==0 else "#FFFFFF"))
                self.set_text_color(*hex_rgb(C["dark"]))
                self.set_font("CJK", "", 8)
                for cell, w in zip(row, widths):
                    self.cell(w, 5, str(cell), border=0, fill=True, align="L")
                self.ln()
            self.ln(1.5)

    def hex_rgb(h):
        h = h.lstrip("#")
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    pdf = PDF()
    img = lambda n: charts.get(n)

    # ====== 封面 ======
    pdf.add_page()
    pdf.set_fill_color(*hex_rgb(C["primary"]))
    pdf.rect(0, 0, 210, 297, "F")
    pdf.set_y(60)
    pdf.set_font("CJK", "B", 34)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 14, "AI 工具部署教程", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("CJK", "", 16)
    pdf.set_text_color(199, 210, 254)
    pdf.cell(0, 9, "从零开始 · 小白也能学会", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 9, "国内用户首选 DeepSeek 方案", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)
    pdf.set_draw_color(129, 140, 248)
    pdf.line(65, pdf.get_y(), 145, pdf.get_y())
    pdf.ln(6)
    pdf.set_font("CJK", "", 11)
    pdf.set_text_color(165, 180, 252)
    pdf.cell(0, 7, "覆盖 Claude Code / Codex / DeepSeek / Cherry Studio", align="C", new_x="LMARGIN", new_y="NEXT")

    # ====== 第1章：认识AI部署 ======
    pdf.add_page()
    pdf.h1("一、认识 AI 工具部署")
    pdf.h2("1.1 什么是 AI 工具部署？")
    pdf.p("AI 工具部署就是把人工智能软件安装到你的电脑上，让它能正常运行。类比：就像给手机装微信、抖音——只是 AI 工具需要一些命令行操作，本教程帮你完成这一步。")
    pdf.h2("1.2 部署后能做什么？")
    pdf.tbl(["场景", "例子"],
            [["编程辅助", "写代码、修Bug、解释逻辑"],
             ["文档处理", "写报告、翻译、润色文章"],
             ["数据分析", "处理Excel、生成图表"],
             ["学习助手", "解答问题、整理笔记"],
             ["自动化", "自动发邮件、整理文件、定时任务"]],
            [40, 130])
    pdf.h2("1.3 国内用户特别说明")
    pdf.tip("本教程提供两条路线：\n• 路线A（海外）：Claude Code + Anthropic 官方 API → 需外网 + 海外支付\n• 路线B（国内推荐）：DeepSeek API + Cherry Studio → 国内直连 + 微信/支付宝\n请根据自己的网络条件选择合适的路线。", border=(29, 86, 219))
    pdf.img(img("deployment_roadmap.png"), 165)

    # ====== 第2章：路线选择 ======
    pdf.add_page()
    pdf.h1("二、选择你的部署路线")
    pdf.h2("2.1 两条路线对比")
    pdf.img(img("ds_vs_claude.png"), 170)
    pdf.h2("2.2 路线推荐")
    pdf.tbl(["条件", "推荐路线", "理由"],
            [["有外网 + Visa卡", "路线A：Claude Code", "Claude 模型能力最强"],
             ["只有国内网", "路线B：DeepSeek + Cherry Studio", "国内直连，微信充值"],
             ["完全零基础", "先用 DeepSeek 网页版", "chat.deepseek.com 免费"],
             ["想都试试", "A+B 双路线", "看哪个更适合你"]],
            [32, 50, 88])
    pdf.tip("建议：无论选择哪条路线，安装步骤大同小异。先按路线B（DeepSeek）走一遍，之后想切换到 Claude 只需改个 API Key。", border=(34, 197, 94))

    # ====== 第3章：准备工作 ======
    pdf.add_page()
    pdf.h1("三、部署前的准备")
    pdf.h2("3.1 环境检查清单")
    pdf.img(img("env"), 170)
    pdf.tip("安装前必读：请逐项核对以上清单。如果某项不满足，先解决再继续。", border=(245, 158, 11))
    pdf.h2("3.2 工具选择")
    pdf.img(img("compare"), 155)
    pdf.tbl(["工具", "合适人群", "难度", "推荐"],
            [["Claude Code CLI", "不介意敲命令", "⭐⭐", "★★★★★"],
             ["Codex CLI", "OpenAI 用户", "⭐⭐", "★★★★★"],
             ["Claude Code 插件", "用 VSCode 的", "⭐", "★★★★★"],
             ["Cherry Studio", "要图形界面", "⭐", "★★★★★"],
             ["OpenClaw", "多功能集成", "⭐⭐⭐", "★★★★"]])

    # ====== 第4章：路线B - DeepSeek 安装（国内推荐）=====
    pdf.add_page()
    pdf.h1("四、路线B：DeepSeek API + Cherry Studio（国内推荐）")
    pdf.h2("4.1 第一步：注册 DeepSeek 并获取 API Key")
    pdf.code(
        "1. 浏览器打开 https://platform.deepseek.com\n"
        "   （国内直接访问，不需要外网！）\n"
        "2. 用手机号注册（+86 即可）\n"
        "3. 登录后 → 左侧「API Keys」→ 创建新 Key\n"
        "4. 复制 Key（sk-...开头）⚠️ 只显示一次！\n"
        "5. 左侧「充值」→ 微信/支付宝扫码 → 最低 ¥10\n"
        "   （¥10 够日常使用 1-2 周，非常便宜）")
    pdf.tip("DeepSeek 价格参考：¥1/百万 token（输入），¥2/百万 token（输出）。\n日常聊天使用一天几分钱，重度使用一个月也就几十块。", border=(34, 197, 94))

    pdf.h2("4.2 第二步：安装 Cherry Studio 客户端")
    pdf.code(
        "1. 浏览器打开 https://cherry-ai.com\n"
        "2. 下载对应版本（Windows .exe / Mac .dmg）\n"
        "3. 双击安装 → 一路下一步\n"
        "4. 启动 Cherry Studio → 左下角「设置」⚙️\n"
        "5. 选择「模型服务」→ 添加 → DeepSeek\n"
        "6. 填入你的 API Key → 保存\n"
        "7. 返回对话界面 → 选择 DeepSeek 模型 → 开始聊天！")
    pdf.tip("Cherry Studio 是国产免费开源客户端，支持 Windows/Mac/Linux，\n界面美观，支持 DeepSeek / OpenAI / 智谱 等多个国内模型。", border=(29, 86, 219))

    pdf.h2("4.3 第三步：验证 & 使用")
    pdf.p("在 Cherry Studio 对话框输入「你好，介绍一下你自己」，看到 DeepSeek 的回复就说明部署成功！")
    pdf.p("Cherry Studio 还支持：文件对话（上传 PDF/Word/Excel 直接提问）、Markdown 预览、多模型同时问答、提示词库等高级功能。")

    # ====== 第5章：路线A - Claude Code 安装 ======
    pdf.add_page()
    pdf.h1("五、路线A：Claude Code 安装（海外用户/有外网）")
    pdf.h2("5.1 总览：四步法")
    pdf.img(img("steps"), 175)

    pdf.h2("5.2 第一步：安装 Node.js")
    pdf.code(
        "浏览器打开 https://nodejs.org → 下载 LTS 版本\n"
        "双击安装，一路 Next，⚠️ 务必勾选「Add to PATH」！\n"
        "验证：终端输入 node --version（应输出 v20.x.x）")
    pdf.tip("看到版本号 = Node.js 安装成功！", border=(34, 197, 94))

    pdf.h2("5.3 第二步：获取 Anthropic API Key")
    pdf.code(
        "1. https://console.anthropic.com（需外网）\n"
        "2. 注册 → API Keys → Create Key → 复制 sk-ant-...\n"
        "3. Billing 页面充值（最低 $5，需 Visa/MasterCard）")
    pdf.tip("安全警告：API Key 是密码，不要发给任何人！", border=(239, 68, 68))

    pdf.h2("5.4 第三步：安装 Claude Code")
    pdf.code(
        "npm install -g @anthropic-ai/claude-code\n"
        "# 等待 2-5 分钟\n"
        "claude --version  # 验证安装")
    pdf.tip("Windows 权限错 → 管理员模式终端（Win+X → 终端(管理员)）", border=(217, 119, 6))

    pdf.h2("5.5 第四步：配置 & 验证")
    pdf.code(
        "claude login  # 粘贴 API Key\n"
        "claude \"你好\"  # 收到回复 = 成功！")

    # ====== Codex CLI 安装 ======
    pdf.h1("五-B、路线A-2：Codex CLI 安装（OpenAI 出品）")
    pdf.h2("什么是 Codex？")
    pdf.p("Codex 是 OpenAI 推出的命令行 AI 编程助手，对标 Claude Code。它同样基于 Node.js，通过 npm 安装，使用 OpenAI API Key。如果你已有 OpenAI 账号（ChatGPT），可以直接用同一个 Key。")
    pdf.h2("安装步骤")
    pdf.code(
        "# 前置：已装 Node.js（见5.2节），已有 OpenAI API Key\n"
        "# 获取 Key：https://platform.openai.com/api-keys\n"
        "# （需外网 + 海外支付，与 Anthropic 类似）\n"
        "\n"
        "npm install -g @openai/codex\n"
        "# 等待安装完成\n"
        "\n"
        "# 配置 Key：\n"
        "# Windows PowerShell:\n"
        '$env:OPENAI_API_KEY = "sk-你的OpenAI Key"\n'
        "# macOS / Linux:\n"
        'export OPENAI_API_KEY="sk-你的OpenAI Key"\n'
        "\n"
        "# 验证：\n"
        "codex \"你好\"\n"
        "# 收到回复 = 部署成功！")
    pdf.tip("Codex vs Claude Code 选择建议：\n• 已有 ChatGPT/OpenAI 账号 → 优先用 Codex\n• 已有 Anthropic 账号 → 优先用 Claude Code\n• 两者功能相似（代码生成、对话、文件操作），选自己已有 Key 的即可\n• 也可两个都装，按需切换", border=(59, 130, 246))

    # ====== Codex + DeepSeek 直驱 ======
    pdf.h1("五-C、最强方案：用 DeepSeek 来跑 Codex")
    pdf.h2("这是什么意思？")
    pdf.p("简单说：Codex 是一个免费的 AI 编程助手（和 Claude Code 差不多），但它本来需要连 OpenAI 才能用。而 DeepSeek 的接口恰好和 OpenAI 是一样的格式——所以只要告诉 Codex「别去 OpenAI，去 DeepSeek 那儿」，Codex 就能用 DeepSeek 的便宜 AI 来工作了。")
    pdf.p("打个比方：Codex 是一辆车，本来只能加 OpenAI 的油。但 DeepSeek 的油枪接口和 OpenAI 一样，所以你只需要把加油站地址从 OpenAI 改成 DeepSeek，车就能跑了。油费只要十分之一。")
    pdf.h2("具体操作（一步步跟着做）")
    pdf.code(
        "# 第一步：去 DeepSeek 拿一个「钥匙」\n"
        "# 用浏览器打开 https://platform.deepseek.com\n"
        "# （这个网站国内直接就能打开，不用翻墙）\n"
        "# 用手机号注册 → 登录 → 点左边「API Keys」→ 点创建\n"
        "# 把生成的钥匙复制下来（sk-...开头的，只显示一次！）\n"
        "# 然后点「充值」，微信扫一下，充 10 块钱就够了\n"
        "\n"
        "# 第二步：装 Codex（在终端里装，第五章已经讲过怎么装）\n"
        "npm install -g @openai/codex\n"
        "\n"
        "# 第三步：告诉 Codex 去 DeepSeek 而不是 OpenAI\n"
        "# 把下面两行复制到终端里运行（先替换成你自己的钥匙）：\n"
        "#\n"
        "# Windows 电脑（PowerShell）：\n"
        '$env:OPENAI_BASE_URL = "https://api.deepseek.com/v1"\n'
        '$env:OPENAI_API_KEY = "sk-粘贴你刚才复制的钥匙"\n'
        "#\n"
        "# Mac 电脑（终端）：\n"
        'export OPENAI_BASE_URL="https://api.deepseek.com/v1"\n'
        'export OPENAI_API_KEY="sk-粘贴你刚才复制的钥匙"\n'
        "\n"
        "# 第四步：试试能不能用\n"
        "codex \"你好\"\n"
        "# 如果它回复你了 → 成功了！Codex 现在用的是 DeepSeek 的 AI\n"
        "# 费用比用 OpenAI 便宜十倍，而且不用翻墙不用外币卡")
    pdf.tip("这三行命令在做什么？\n• 第一行：告诉 Codex「换地址，去 DeepSeek 别去 OpenAI」\n• 第二行：把你的 DeepSeek 钥匙交给 Codex\n• 第三行：测试一下能不能正常聊天\n\n以后每次打开终端要先用 codex 的时候，都要先输前两行。如果想一劳永逸，看附录 9.3 的环境变量永久配置。", border=(34, 197, 94))

    pdf.h2("方案二：Claude Code 能不能也这样？")
    pdf.p("Claude Code 稍微麻烦一点——它用的是另一种接口格式，不能直接连 DeepSeek。需要一个「转接头」（也就是下一章要讲的中转站），把 Claude 的话翻译成 DeepSeek 能听懂的格式。具体看第六章。")

    # ====== 第6章：中转站方案 ======
    pdf.add_page()
    pdf.h1("六、进阶：让 Claude Code 也能用 DeepSeek")

    pdf.h2("6.1 为什么需要「中转站」？")
    pdf.p("Claude Code 说的是「英语」，DeepSeek 说的是「中文」，它们不能直接对话。中转站就像一个翻译，Claude Code 把请求发给中转站，中转站翻译成 DeepSeek 能懂的格式，DeepSeek 回答后再翻译回来。整个过程在国内网络完成，不需要翻墙。")
    pdf.p("你需要做的就是：找一个靠谱的中转站 → 注册拿到它的地址和钥匙 → 告诉 Claude Code 去那里。")

    pdf.h2("6.2 具体操作")
    pdf.code(
        "# 第一步：找一个中转站（提供 API 中转服务的网站）\n"
        "# 选站建议：支持微信/支付宝充值、运营超过半年、有客服群\n"
        "# 注册后你会拿到两样东西：\n"
        "#   ① 中转站地址（一串网址，类似 https://xxx.com/v1）\n"
        "#   ② API 钥匙（类似 sk-... 的字符串）\n"
        "\n"
        "# 第二步：告诉 Claude Code 用中转站\n"
        "# 把下面两行复制到终端（替换成你自己的地址和钥匙）：\n"
        "#\n"
        "# Windows 电脑（PowerShell）：\n"
        '$env:ANTHROPIC_BASE_URL = "粘贴中转站给你的地址"\n'
        '$env:ANTHROPIC_API_KEY = "粘贴中转站给你的钥匙"\n'
        "#\n"
        "# Mac 电脑（终端）：\n"
        'export ANTHROPIC_BASE_URL="粘贴中转站给你的地址"\n'
        'export ANTHROPIC_API_KEY="粘贴中转站给你的钥匙"\n'
        "\n"
        "# 第三步：正常使用\n"
        "claude \"你好\"\n"
        "# Claude Code 现在通过中转站走 DeepSeek 了！")
    pdf.tip("中转站避坑提醒：\n• 先用小金额（充 5-10 块）测试几天，稳定了再多充\n• 不要相信「白菜价」「不限量」的宣传，正常价格是几块钱/百万token\n• 优先选有 QQ 群或微信群的，出了问题能找到人\n• 建议同时注册 2 个中转站备用，一个挂了还有另一个", border=(245, 158, 11))

    pdf.h2("6.3 全部方案总结")
    pdf.tbl(["方案", "门槛", "费用", "推荐场景"],
            [["DeepSeek + Cherry Studio", "国内网 + 微信", "¥10起", "⭐ 小白首选（图形界面）"],
             ["Codex + DeepSeek(直连)", "国内网 + 微信", "¥10起", "⭐ 编程首选（命令行）"],
             ["Codex + OpenAI", "外网 + Visa", "$5起", "已有ChatGPT账号"],
             ["Claude Code + Anthropic", "外网 + Visa", "$5起", "需最强模型"],
             ["Claude Code + 中转站", "国内网 + 微信", "¥10起", "想用CC但没外网"],
             ["DeepSeek 网页版", "浏览器就行", "免费", "先体验再决定"]])

    # ====== 第7章：问题排查 ======
    pdf.add_page()
    pdf.h1("七、常见问题排查")
    pdf.img(img("trouble"), 170)
    pdf.h2("7.1 通用问题")
    for title, content in [
        ("npm: command not found", "Node.js 没装好 → 重装 nodejs.org LTS版，勾选 Add to PATH"),
        ("权限错误 EACCES", "Mac/Linux 加 sudo，Windows 用管理员终端"),
        ("claude: command not found", "npm list -g 确认已安装，重启终端，或用 npx @anthropic-ai/claude-code"),
        ("API 401 错误", "Key 无效/过期 → 去对应平台重新生成；检查账户余额"),
        ("网络超时", "关 VPN 试，或切手机热点排除宽带问题"),
        ("DeepSeek 连不上", "确认地址是 platform.deepseek.com，试试用手机流量访问"),
    ]:
        pdf.h3(title)
        pdf.tip(content, border=(239, 68, 68) if "401" in title or "not found" in title else (217, 119, 6))

    # ====== 第8章：使用技巧 ======
    pdf.add_page()
    pdf.h1("八、安装后使用指南")
    pdf.h2("8.1 常用命令速查")
    pdf.tbl(["命令", "作用"],
            [["claude / cherry studio 对话框", "交互式对话"],
             ['claude "问题"', "单次提问"],
             ["多轮追问", "不满意就让它改"],
             ["上传文件", "PDF/Excel/Word 直接拖入"],
             ["claude --version", "查看版本"]],
            [70, 100])
    pdf.h2("8.2 高效技巧")
    pdf.tbl(["技巧", "说明"],
            [["越具体越好", "模糊问题→模糊答案，给细节"],
             ["给上下文", "告诉AI你的背景和目的"],
             ["可以追问", "太长→压缩，不对→换个角度"],
             ["定期更新", "npm update -g @anthropic-ai/claude-code"]])

    pdf.h2("8.3 费用参考")
    pdf.tbl(["服务", "价格", "说明"],
            [["DeepSeek API", "¥1-2/百万token", "极便宜，国内直连"],
             ["OpenAI API", "$0.15-15/百万token", "Codex 依赖"],
             ["Anthropic API", "$3-15/百万token", "贵但模型最强"],
             ["Cherry Studio", "免费", "开源客户端"],
             ["Codex / Claude Code CLI", "免费", "开源命令行工具"]],
            [42, 42, 86])

    # ====== 附录 ======
    pdf.add_page()
    pdf.h1("九、附录")
    pdf.h2("9.1 速查表")
    pdf.tbl(["步骤", "路线A (Claude)", "路线A (Codex)", "路线B (DeepSeek)"],
            [["1", "装 Node.js", "装 Node.js", "注册 DeepSeek"],
             ["2", "注册 Anthropic(需外网)", "注册 OpenAI(需外网)", "微信/支付宝充值"],
             ["3", "npm install -g claude-code", "npm install -g @openai/codex", "下载 Cherry Studio"],
             ["4", "claude login 配 Key", "配置 OPENAI_API_KEY", "填入 Key → 开始用"]],
            [10, 52, 52, 56])

    pdf.h2("9.2 重要链接")
    pdf.tbl(["资源", "链接"],
            [["DeepSeek 控制台", "https://platform.deepseek.com"],
             ["Cherry Studio 下载", "https://cherry-ai.com"],
             ["OpenAI 控制台", "https://platform.openai.com"],
             ["Anthropic 控制台", "https://console.anthropic.com"],
             ["Node.js 下载", "https://nodejs.org"],
             ["Claude Code 文档", "https://docs.anthropic.com/en/docs/claude-code"],
             ["Codex CLI 文档", "https://github.com/openai/codex"]],
            [52, 118])

    pdf_path = OUTPUT_DIR / "AI工具部署教程_小白版.pdf"
    try:
        pdf.output(str(pdf_path))
    except PermissionError:
        pdf_path = OUTPUT_DIR / "AI工具部署教程_小白版_v2.pdf"
        pdf.output(str(pdf_path))
        print(f"  ⚠️ 原文件被占用，输出到: {pdf_path.name}")
    print(f"\n✅ PDF: {pdf_path} ({pdf_path.stat().st_size//1024} KB)")
    return pdf_path


# ============================================================
# PPT 生成（python-pptx，15页）
# ============================================================
def build_ppt_old(charts):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    # 使用用户的PPT模板（继承主题色/字体/背景样式）
    template_path = Path("D:/ProjectManagement/Claude/ppt模板.pptx")
    if template_path.exists():
        prs = Presentation(str(template_path))
        # 删除模板中所有现有幻灯片（保留主题/母版）
        sldIdLst = prs._element.sldIdLst
        if sldIdLst is not None:
            for sldId in list(sldIdLst):
                rId = sldId.get(qn('r:id'))
                sldIdLst.remove(sldId)
                if rId is not None:
                    try:
                        prs.part.drop_rel(rId)
                    except Exception:
                        pass
        # 使用第一个可用布局（如果有blank布局优先用）
        blank = prs.slide_layouts[-1] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        print(f"  ✅ 使用模板: {template_path.name} (主题已继承)")
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        print(f"  ⚠️ 模板未找到，使用默认主题")

    def hex_rgb(h):
        h = h.lstrip("#")
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    def bg(s, c):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(*hex_rgb(c))

    def tb(s, l, t, w, h, txt, sz=14, c=C["dark"], b=False, a=PP_ALIGN.LEFT):
        box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = txt
        p.font.size = Pt(sz); p.font.color.rgb = RGBColor(*hex_rgb(c))
        p.font.bold = b; p.font.name = "Microsoft YaHei"; p.alignment = a
        return tf

    def card(s, l, t, w, h, color):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(*hex_rgb(color)); sh.line.fill.background()
        return sh

    def img(s, path, l, t, w, h=None):
        p = Path(path) if isinstance(path, str) else path
        if p and p.exists():
            kw = {"width": Inches(w)}
            if h: kw["height"] = Inches(h)
            s.shapes.add_picture(str(p), Inches(l), Inches(t), **kw)

    im = lambda n: charts.get(n)

    # === S1: 封面 ===
    s = prs.slides.add_slide(blank); bg(s, C["primary"])
    tb(s, 1, 1.5, 11.3, 1.5, "AI 工具部署教程", 48, "#FFFFFF", True, PP_ALIGN.CENTER)
    tb(s, 1, 3.2, 11.3, 1, "从零开始 · 小白也能学会", 22, "#C7D2FE", a=PP_ALIGN.CENTER)
    tb(s, 1, 3.8, 11.3, 0.8, "国内用户首选 DeepSeek + Cherry Studio", 18, "#A5B4FC", a=PP_ALIGN.CENTER)
    tb(s, 1, 5.5, 11.3, 0.5, "Claude Code / Codex / DeepSeek / Cherry Studio 全覆盖", 13, "#818CF8", a=PP_ALIGN.CENTER)

    # === S2: 目录 ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 8, 0.6, "课程目录", 30, C["primary"], True)
    toc = [("01","认识 AI 工具部署"), ("02","选择部署路线"), ("03","部署前准备"),
           ("04","路线B: DeepSeek(推荐)"), ("05","路线A: Claude Code"), ("06","路线A: Codex CLI"),
           ("07","DeepSeek驱动Codex/CC"), ("08","进阶: 中转站方案"),
           ("09","常见问题排查"), ("10","使用指南"), ("11","附录")]
    for i, (n, t) in enumerate(toc):
        y = 1.3 + i * 0.65
        c = card(s, 1.2, y, 0.5, 0.5, C["primary"])
        c.text_frame.paragraphs[0].text = n
        c.text_frame.paragraphs[0].font.size = Pt(14)
        c.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(s, 2.0, y, 8, 0.4, t, 18, C["dark"], True)

    # === S3: 路线对比 ===
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC")
    tb(s, 0.5, 0.1, 8, 0.5, "两条部署路线对比", 28, C["primary"], True)
    img(s, im("ds_vs_claude"), 0.5, 0.8, 12.3)

    # === S4: 环境检查 & 工具对比（缩小图片防重叠）===
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC")
    tb(s, 0.5, 0.05, 8, 0.4, "环境检查 & 工具选择", 26, C["primary"], True)
    img(s, im("env"), 0.3, 0.5, 6.0)
    img(s, im("compare"), 6.8, 0.5, 6.0)

    # === S5: 路线B-DeepSeek注册 ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hex_rgb(C["ds"])); bar.line.fill.background()
    tb(s, 0.6, 1.5, 3.5, 1, "路线B", 22, "#FFFFFF", True)
    tb(s, 0.6, 2.5, 3.5, 1.5, "DeepSeek API\n+ Cherry Studio", 34, "#FFFFFF", True)
    tb(s, 0.6, 4.5, 3.5, 1.5, "🇨🇳 国内直连\n💰 微信/支付宝充值\n⏱️ 约10分钟完成", 14, "#BFDBFE")
    items = ["① platform.deepseek.com 注册", "② 手机号注册 (+86)", "③ API Keys → 创建 Key → 复制",
             "④ 微信/支付宝充值 ¥10", "⑤ cherry-ai.com 下载客户端", "⑥ 填入 Key → 开始使用！"]
    for j, item in enumerate(items):
        y = 0.8 + j * 0.95
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(y+0.12), Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(*hex_rgb(C["ds"])); dot.line.fill.background()
        tb(s, 5.6, y, 7, 0.55, item, 14, C["dark"])

    # === S6: Cherry Studio 介绍 ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 10, 0.6, "Cherry Studio — 国产免费 AI 客户端", 28, C["ds"], True)
    features = [
        ("🌐 国内直连", "无需外网，下载即用"),
        ("🔌 多模型支持", "DeepSeek / OpenAI / 智谱 / 通义千问..."),
        ("📂 文件对话", "拖入 PDF/Word/Excel 直接提问"),
        ("🎨 美观界面", "现代化设计，比命令行友好"),
        ("💰 完全免费", "开源软件，Apache 2.0 协议"),
        ("🔧 高级功能", "多模型对比、提示词库、Markdown 渲染"),
    ]
    for j, (title, desc) in enumerate(features):
        x = 0.5 + (j % 3) * 4.1
        y = 1.2 + (j // 3) * 2.9
        card(s, x, y, 3.6, 2.5, "#F0F7FF")
        tb(s, x+0.15, y+0.15, 3.3, 0.45, title, 17, C["ds"], True)
        tb(s, x+0.15, y+0.9, 3.3, 0.9, desc, 12, C["dark"])

    # === S7: 路线A-四步法 ===
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC")
    tb(s, 0.5, 0.1, 8, 0.5, "路线A：Claude Code 四步安装法", 28, C["primary"], True)
    img(s, im("steps"), 0.5, 0.8, 12.3)

    # === S8: Claude Code 详细 ===
    for idx, (st, title, color, items) in enumerate([
        ("STEP 1+2", "环境 + API Key", C["success"],
         ["装 Node.js (nodejs.org → LTS)", "⚠️ 勾选 Add to PATH",
          "注册 Anthropic (需外网)", "创建 API Key → 保存好"]),
        ("STEP 3+4", "安装 + 验证", C["primary"],
         ["npm install -g @anthropic-ai/claude-code", "等待 2-5 分钟",
          "claude login → 粘贴 Key", 'claude "你好" → 成功！']),
    ]):
        s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hex_rgb(color)); bar.line.fill.background()
        tb(s, 0.8, 2, 3.2, 1, st, 22, "#FFFFFF", True)
        tb(s, 0.8, 3, 3.2, 1.5, title, 32, "#FFFFFF", True)
        for j, item in enumerate(items):
            y = 1.2 + j * 1.1
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(y+0.12), Inches(0.15), Inches(0.15))
            dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(*hex_rgb(color)); dot.line.fill.background()
            tb(s, 5.6, y, 7, 0.55, item, 14, C["dark"], "⚠️" in item)

    # === S9: Codex CLI ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hex_rgb("#10A37F")); bar.line.fill.background()
    tb(s, 0.6, 1.5, 3.5, 1, "路线A-2", 22, "#FFFFFF", True)
    tb(s, 0.6, 2.5, 3.5, 1.5, "Codex CLI\nOpenAI 出品", 34, "#FFFFFF", True)
    tb(s, 0.6, 4.5, 3.5, 1.5, "🤖 OpenAI 官方 CLI 工具\n🔑 可用 ChatGPT 同款 Key\n⚡ npm 一键安装", 14, "#D1FAE5")
    items = ["① platform.openai.com 获取 Key", "② npm install -g @openai/codex", "③ 配置环境变量 OPENAI_API_KEY",
             "④ codex \"你好\" 验证", "⑤ 与 Claude Code 功能相似", "⑥ 可按需同时安装两者"]
    for j, item in enumerate(items):
        y = 0.8 + j * 0.95
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(y+0.12), Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(*hex_rgb("#10A37F")); dot.line.fill.background()
        tb(s, 5.6, y, 7, 0.55, item, 14, C["dark"])

    # === S10: DeepSeek API 直驱方案（白话版）===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hex_rgb("#059669")); bar.line.fill.background()
    tb(s, 0.6, 1.0, 3.5, 1, "最强方案", 22, "#FFFFFF", True)
    tb(s, 0.6, 2.0, 3.5, 2, "用 DeepSeek\n来跑 Codex", 30, "#FFFFFF", True)
    tb(s, 0.6, 4.2, 3.5, 2.5, "Codex 是车，DeepSeek 是油\n接口一样，换个地址就能用\n费用只要 OpenAI 的 1/10\n国内直连 + 微信充值", 12, "#D1FAE5")
    items = ["去 DeepSeek 拿钥匙: platform.deepseek.com",
             "装 Codex: npm install -g @openai/codex",
             "告诉 Codex 去 DeepSeek（复制粘贴一行）",
             "第一行: OPENAI_BASE_URL=api.deepseek.com/v1",
             "第二行: OPENAI_API_KEY=你的钥匙",
             "试试: codex 你好 → 用 DeepSeek 回复!"]
    for j, item in enumerate(items):
        y = 0.6 + j * 0.95
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(y+0.12), Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(*hex_rgb("#059669")); dot.line.fill.background()
        is_code = "BASE_URL" in item or "API_KEY" in item
        tb(s, 5.6, y, 7, 0.55, item, 12 if is_code else 13, "#1E293B", is_code)

    # === S11: 中转站方案（白话版）===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    tb(s, 0.8, 0.2, 11, 0.5, "进阶：让 Claude Code 走 DeepSeek（中转站）", 20, C["accent"], True)
    tb(s, 0.8, 0.8, 11.5, 0.5, "Claude Code 说英语，DeepSeek 说中文，中转站当翻译——在国内就能完成。", 13, C["dark"])
    card(s, 0.8, 1.5, 11.5, 3.2, "#F5F3FF")
    tb(s, 1.2, 1.65, 10.5, 0.3, "只需三步：", 14, C["accent"], True)
    tb(s, 1.2, 2.1, 10.5, 2.3,
        '1. 找中转站 → 微信注册 → 拿到地址和钥匙\n\n'
        '2. 打开终端，复制粘贴两行（替换成你自己的）：\n'
        '   $env:ANTHROPIC_BASE_URL = "中转站给的地址"\n'
        '   $env:ANTHROPIC_API_KEY = "中转站给的钥匙"\n\n'
        '3. 正常用: claude "你好" → 走DeepSeek回复!',
        10, "#334155")
    tb(s, 0.8, 5.0, 11.5, 0.35, "选站提醒：充5-10块先试几天。别信「白菜价」宣传，认准运营半年以上的站。", 11, C["warning"], True)
    card(s, 0.8, 5.5, 11.5, 1.7, "#F8FAFC")
    tb(s, 1.2, 5.65, 10.5, 0.3, "方案速查", 13, C["dark"], True)
    for j, (plan, cost, scene) in enumerate([
        ("DS+Cherry", "¥10起", "小白首选"),
        ("Codex+DS", "¥10起", "编程首选"),
        ("Codex+OAI", "$5起", "有ChatGPT"),
        ("CC+中转", "¥10起", "没外网"),
    ]):
        x = 1.2 + j * 2.95
        tb(s, x, 6.1, 2.8, 0.3, f"{plan}", 9, C["dark"], True)
        tb(s, x, 6.45, 2.8, 0.25, f"{cost} | {scene}", 8, "#64748B")

    # === S10: 问题排查 ===
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC")
    tb(s, 0.5, 0.1, 8, 0.5, "常见问题排查", 28, C["primary"], True)
    img(s, im("trouble"), 0.3, 0.7, 12.7)

    # === S11: 使用技巧 ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 10, 0.6, "高效使用技巧", 28, C["primary"], True)
    tips = [("越具体越好","❌写方案 → ✅写奶茶店运营方案"), ("给上下文","告诉AI你的背景和目的"),
            ("可以追问","太长→压缩 / 不对→换个角度"), ("文件对话","拖入PDF/Excel直接提问"),
            ("定期更新","npm update 或客户端自动更新")]
    for j, (t, d) in enumerate(tips):
        x = 0.5 + j * 2.5
        card(s, x, 1.3, 2.2, 5.3, "#F8FAFC")
        tb(s, x+0.15, 1.45, 1.9, 0.45, t, 15, C["dark"], True)
        tb(s, x+0.15, 2.2, 1.9, 2.0, d, 11, "#64748B")

    # === S12: 命令速查 ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 8, 0.6, "命令 & 操作速查", 28, C["primary"], True)
    cmds = [("Claude Code 交互", "终端输入 claude"), ("Codex 交互", "终端输入 codex"),
            ("查看版本", "claude/codex --version"), ("重新登录", "claude login"),
            ("Cherry Studio", "图形界面 → 选模型 → 对话"), ("更新工具", "npm update -g @anthropic-ai/claude-code")]
    for j, (cmd, desc) in enumerate(cmds):
        y = 1.3 + j * 0.85
        card(s, 1, y, 5.5, 0.65, "#F1F5F9")
        tb(s, 1.2, y+0.1, 5.2, 0.45, cmd, 15, C["primary"], True)
        tb(s, 7, y+0.1, 5.5, 0.45, desc, 14, "#475569")

    # === S13: 费用对比 ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 8, 0.6, "费用对比", 28, C["primary"], True)
    fees = [("DeepSeek API","极便宜","¥1-2/百万token · 国内直连 · 微信充值", C["success"]),
            ("OpenAI API","中等","$0.15-15/百万token · 需外网 · Visa", C["info"]),
            ("Anthropic API","较贵","$3-15/百万token · 需外网 · Visa", C["warning"]),
            ("Cherry/Codex/CC","免费","开源客户端/CLI · 永久免费", C["primary"])]
    for j, (name, price, note, col) in enumerate(fees):
        y = 1.3 + j * 1.5
        card(s, 1, y, 11, 1.2, "#F8FAFC")
        bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(0.07), Inches(1.2))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = RGBColor(*hex_rgb(col)); bar2.line.fill.background()
        tb(s, 1.5, y+0.1, 3, 0.4, name, 19, C["dark"], True)
        tb(s, 5.5, y+0.1, 2.5, 0.4, price, 20, col, True)
        tb(s, 1.5, y+0.6, 9, 0.35, note, 13, "#64748B")

    # === S14: 重要链接 ===
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 8, 0.6, "重要链接汇总", 28, C["primary"], True)
    links = [("DeepSeek 控制台","platform.deepseek.com","注册 + 获取 Key + 充值"),
             ("Cherry Studio","cherry-ai.com","下载免费 AI 客户端"),
             ("Node.js","nodejs.org","JavaScript 运行环境"),
             ("Anthropic 控制台","console.anthropic.com","Claude API Key（需外网）"),
             ("Claude Code 文档","docs.anthropic.com/en/docs/claude-code","官方文档")]
    for j, (name, url, desc) in enumerate(links):
        y = 1.4 + j * 1.1
        card(s, 1, y, 11, 0.9, "#F8FAFC")
        tb(s, 1.3, y+0.1, 3, 0.35, name, 17, C["dark"], True)
        tb(s, 5, y+0.1, 4.5, 0.35, url, 15, C["primary"])
        tb(s, 1.3, y+0.5, 8, 0.3, desc, 12, "#64748B")

    # === S15: 结束页 ===
    s = prs.slides.add_slide(blank); bg(s, C["primary"])
    tb(s, 1, 1.8, 11.3, 1, "恭喜完成学习！", 42, "#FFFFFF", True, PP_ALIGN.CENTER)
    tb(s, 1, 3.0, 11.3, 1, "现在你已经可以独立完成 AI 工具的部署", 20, "#C7D2FE", a=PP_ALIGN.CENTER)
    nexts = [("小白首选", "DeepSeek + Cherry Studio"), ("代码专家", "Codex + DeepSeek"), ("全能王者", "Claude Code + 中转站")]
    for j, (t, d) in enumerate(nexts):
        x = 2 + j * 3.3
        card(s, x, 4.2, 2.8, 1.8, "#6366F1")
        tb(s, x+0.2, 4.4, 2.4, 0.5, t, 16, "#FFFFFF", True)
        tb(s, x+0.2, 5.1, 2.4, 0.5, d, 12, "#C7D2FE")
    tb(s, 1, 6.5, 11.3, 0.5, "让每个人都能用上 AI", 12, "#818CF8", a=PP_ALIGN.CENTER)

    ppt_path = OUTPUT_DIR / "AI工具部署教程_小白版.pptx"
    try:
        prs.save(str(ppt_path))
    except PermissionError:
        ppt_path = OUTPUT_DIR / "AI工具部署教程_小白版_v2.pptx"
        prs.save(str(ppt_path))
        print(f"  ⚠️ 原文件被占用，输出到: {ppt_path.name}")
    print(f"✅ PPT: {ppt_path} ({ppt_path.stat().st_size//1024} KB, {len(prs.slides)} slides)")
    return ppt_path


# ============================================================
# Excel 生成（openpyxl，多Sheet表格版）
# ============================================================
def build_excel():
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # 通用样式
    hdr_font = Font(name="Microsoft YaHei", size=11, bold=True, color="FFFFFF")
    hdr_fill = PatternFill(start_color="4F46E5", end_color="4F46E5", fill_type="solid")
    hdr_fill_green = PatternFill(start_color="059669", end_color="059669", fill_type="solid")
    hdr_fill_blue = PatternFill(start_color="1A56DB", end_color="1A56DB", fill_type="solid")
    body_font = Font(name="Microsoft YaHei", size=10)
    title_font = Font(name="Microsoft YaHei", size=14, bold=True, color="4F46E5")
    wrap_align = Alignment(wrap_text=True, vertical="center")
    center_align = Alignment(horizontal="center", vertical="center")
    thin_border = Border(
        left=Side(style="thin", color="E2E8F0"),
        right=Side(style="thin", color="E2E8F0"),
        top=Side(style="thin", color="E2E8F0"),
        bottom=Side(style="thin", color="E2E8F0"),
    )

    def add_title(ws, text, ncols=4):
        """在第一行写入大标题（跨多列展示在首列）"""
        c = ws.cell(row=1, column=1, value=text)
        c.font = title_font; c.alignment = Alignment(vertical="center")
        ws.row_dimensions[1].height = 28

    def style_sheet(ws, headers, rows, col_widths, fill="primary", start_row=2):
        fills = {"primary": hdr_fill, "green": hdr_fill_green, "blue": hdr_fill_blue}
        hf = fills.get(fill, hdr_fill)
        for ci, h in enumerate(headers, 1):
            c = ws.cell(row=start_row, column=ci, value=h)
            c.font = hdr_font; c.fill = hf; c.alignment = center_align; c.border = thin_border
        for ri, row in enumerate(rows, start_row + 1):
            bg = "F8FAFC" if ri % 2 == 0 else "FFFFFF"
            row_fill = PatternFill(start_color=bg, end_color=bg, fill_type="solid")
            for ci, val in enumerate(row, 1):
                c = ws.cell(row=ri, column=ci, value=val)
                c.font = body_font; c.fill = row_fill; c.alignment = wrap_align; c.border = thin_border
        for ci, w in enumerate(col_widths, 1):
            ws.column_dimensions[get_column_letter(ci)].width = w
        ws.freeze_panes = f"A{start_row + 1}"

    # === Sheet 1: 总览 ===
    ws = wb.active; ws.title = "总览"
    add_title(ws, "AI 工具部署教程 — 小白版")
    overview = [
        ["一、认识 AI 工具部署", "把 AI 软件装到电脑上，让它正常运行", "类比给手机装 App"],
        ["二、三条部署路线", "A: Claude Code/Codex(海外) | B: DeepSeek+Cherry Studio(国内)", "选适合自己的"],
        ["三、部署前准备", "OS/Win10+或Mac12+、8G内存、2G硬盘、稳定网络", "逐项核对检查清单"],
        ["四、DeepSeek 路线", "注册 platform.deepseek.com → 微信充值 → Cherry Studio", "⭐ 小白首选"],
        ["五、Claude Code 路线", "Node.js → Anthropic Key → npm install → claude login", "需外网+Visa"],
        ["五-B、Codex CLI 路线", "Node.js → OpenAI Key → npm install → 配置环境变量", "需外网+Visa"],
        ["六、中转站方案", "设置 ANTHROPIC_BASE_URL → 走国内API", "进阶用户"],
        ["七、常见问题", "npm not found / 权限错 / 401 / 网络超时", "本文档有详细解决方案"],
        ["八、使用技巧", "越具体越好 / 给上下文 / 可追问 / 文件对话", "提升效率"],
        ["九、附录", "速查表 / 重要链接 / 环境变量配置", "快速查阅"],
    ]
    style_sheet(ws, ["章节", "核心内容", "备注"], overview, [25, 52, 28])

    # === Sheet 2: 路线对比 ===
    ws = wb.create_sheet("路线对比")
    add_title(ws, "部署路线对比一览")
    routes = [
        ["DeepSeek + Cherry Studio", "国内网 + 微信", "¥10起", "⭐ 小白首选", "10分钟"],
        ["Codex + DeepSeek(直驱)", "国内网 + 微信", "¥10起", "⭐ 编程首选", "15分钟"],
        ["Codex + OpenAI", "外网 + Visa", "$5起", "已有ChatGPT账号", "15分钟"],
        ["Claude Code + Anthropic", "外网 + Visa", "$5起", "需最强模型", "15分钟"],
        ["Claude Code + 中转站", "国内网 + 微信", "¥10起", "想用CC但没外网", "15分钟"],
        ["DeepSeek 网页版", "浏览器就行", "免费", "先体验再决定", "0分钟"],
    ]
    style_sheet(ws, ["方案", "门槛", "起步费用", "推荐场景", "预计时间"], routes, [28, 22, 15, 26, 15], "blue")

    # === Sheet 3: DeepSeek安装 ===
    ws = wb.create_sheet("DeepSeek安装")
    add_title(ws, "路线B：DeepSeek API + Cherry Studio（国内推荐）")
    ds_rows = [
        ["Step 1", "注册 DeepSeek", "platform.deepseek.com → 手机号注册 → API Keys → 创建Key"],
        ["Step 2", "充值", "微信/支付宝扫码 → 最低 ¥10"],
        ["Step 3", "下载 Cherry Studio", "cherry-ai.com → 下载对应版本 → 安装"],
        ["Step 4", "配置", "Cherry Studio → 设置 → 模型服务 → 添加 DeepSeek → 填入Key"],
        ["Step 5", "验证", "对话框输入「你好」→ 收到回复 = 成功"],
    ]
    style_sheet(ws, ["步骤", "操作", "详细说明"], ds_rows, [10, 22, 78], "green")

    # === Sheet 4: Claude Code安装 ===
    ws = wb.create_sheet("Claude Code安装")
    add_title(ws, "路线A-1：Claude Code 安装（Anthropic）")
    cc_rows = [
        ["Step 1", "安装 Node.js", "nodejs.org → 下载LTS → 安装(勾选Add to PATH)", "node --version"],
        ["Step 2", "获取 API Key", "console.anthropic.com(需外网) → API Keys → Create", "sk-ant-..."],
        ["Step 3", "安装 Claude Code", "npm install -g @anthropic-ai/claude-code", "claude --version"],
        ["Step 4", "配置", "claude login → 粘贴 Key", "看到 ✓ 即成功"],
    ]
    style_sheet(ws, ["步骤", "做什么", "详细说明", "验证命令"], cc_rows, [10, 22, 62, 22])

    # === Sheet 5: Codex安装 ===
    ws = wb.create_sheet("Codex安装")
    add_title(ws, "路线A-2：Codex CLI 安装（OpenAI）")
    cx_rows = [
        ["Step 1", "安装 Node.js", "nodejs.org → 下载LTS → 安装(勾选Add to PATH)", "node --version"],
        ["Step 2", "获取 OpenAI Key", "platform.openai.com(需外网) → API Keys → Create", "sk-proj-..."],
        ["Step 3", "安装 Codex", "npm install -g @openai/codex", "codex --version"],
        ["Step 4", "配置环境变量", "Win: $env:OPENAI_API_KEY='sk-...'", 'codex "你好"'],
        ["Step 5", "验证", "终端输入 codex \"你好\" → 收到回复 = 成功", "开始使用!"],
    ]
    style_sheet(ws, ["步骤", "做什么", "详细说明", "验证/命令"], cx_rows, [10, 22, 62, 22])

    # === Sheet 6: DeepSeek直驱Codex ===
    ws = wb.create_sheet("DeepSeek直驱Codex")
    add_title(ws, "最强方案：DeepSeek API 驱动 Codex CLI（国内直连 + 超低价）")
    dd_rows = [
        ["原理", "DeepSeek API 与 OpenAI 格式完全兼容", "Codex 可直接连 DeepSeek，无需 OpenAI 账号"],
        ["Step 1", "获取 DeepSeek Key", "platform.deepseek.com → API Keys → 微信充值 ¥10"],
        ["Step 2", "安装 Codex", "npm install -g @openai/codex"],
        ["Step 3", "设置环境变量", 'OPENAI_BASE_URL=https://api.deepseek.com/v1'],
        ["Step 3", "设置环境变量", 'OPENAI_API_KEY=sk-你的DeepSeek_Key'],
        ["Step 4", "验证", 'codex "你好" → 通过 DeepSeek 运行！'],
        ["费用", "Codex CLI: 免费", "DeepSeek API: ¥1/百万token（OpenAI的1/10）"],
        ["切换回OpenAI", "去掉 BASE_URL 即可", "或设回 https://api.openai.com/v1"],
    ]
    style_sheet(ws, ["项目", "操作", "详细说明"], dd_rows, [18, 30, 78], "green")

    # === Sheet 7: 问题排查 ===
    ws = wb.create_sheet("问题排查")
    add_title(ws, "常见问题 & 解决方案")
    trouble = [
        ["npm: command not found", "Node.js 没装好", "重装 nodejs.org LTS版，勾选 Add to PATH，重启终端"],
        ["权限错误 EACCES", "没有系统写入权限", "Mac/Linux: 加 sudo | Windows: 管理员终端"],
        ["claude/codex: not found", "全局安装路径不在PATH", "npm list -g 确认安装，重启终端，或用 npx 代替"],
        ["API 401 错误", "Key 无效/过期/欠费", "去对应平台重新生成Key，检查账户余额"],
        ["网络超时/连不上", "网络/防火墙问题", "关VPN试，切手机热点，检查防火墙"],
        ["DeepSeek 连不上", "地址或网络问题", "确认 platform.deepseek.com，用手机流量试"],
    ]
    style_sheet(ws, ["报错信息", "原因", "解决方案"], trouble, [30, 24, 90], "primary")

    # === Sheet 7: 命令速查 ===
    ws = wb.create_sheet("命令速查")
    add_title(ws, "常用命令 & 操作速查")
    cmds = [
        ["claude", "进入 Claude Code 交互模式", "终端直接输入"],
        ["codex", "进入 Codex CLI 交互模式", "终端直接输入"],
        ['claude "问题"', "Claude Code 单次提问", "不进入交互模式"],
        ['codex "问题"', "Codex CLI 单次提问", "不进入交互模式"],
        ["claude/codex --version", "查看版本号", "确认安装成功"],
        ["claude login", "重新登录/更换 API Key", "Key 过期时使用"],
        ["npm update -g @anthropic-ai/claude-code", "更新 Claude Code", "建议每周执行"],
        ["npm update -g @openai/codex", "更新 Codex CLI", "建议每周执行"],
        ["npm uninstall -g @anthropic-ai/claude-code", "卸载 Claude Code", "不再使用时"],
        ["Cherry Studio 操作", "图形界面选模型 → 对话", "支持文件上传提问"],
    ]
    style_sheet(ws, ["命令/操作", "作用", "备注"], cmds, [30, 42, 42])

    # === Sheet 8: 费用参考 ===
    ws = wb.create_sheet("费用参考")
    add_title(ws, "各服务费用对比")
    fees = [
        ["DeepSeek API", "¥1-2 / 百万 token", "极便宜，国内直连", "微信/支付宝"],
        ["Codex + DeepSeek", "¥1-2 / 百万 token", "Codex免费+DeepSeek低价", "微信/支付宝"],
        ["OpenAI API (GPT-4o)", "$2.5-10 / 百万 token", "中等", "Visa/MasterCard"],
        ["OpenAI API (GPT-4o-mini)", "$0.15-0.6 / 百万 token", "很便宜", "Visa/MasterCard"],
        ["Anthropic API (Sonnet)", "$3 / 百万 token", "中等偏贵", "Visa/MasterCard"],
        ["Anthropic API (Opus)", "$15 / 百万 token", "最贵", "Visa/MasterCard"],
        ["Cherry Studio 客户端", "免费", "开源软件", "无需付费"],
        ["Claude Code CLI", "免费", "开源命令行工具", "无需付费"],
        ["Codex CLI", "免费", "开源命令行工具", "无需付费"],
    ]
    style_sheet(ws, ["服务/工具", "价格", "评价", "支付方式"], fees, [30, 28, 22, 22])

    # === Sheet 9: 重要链接 ===
    ws = wb.create_sheet("重要链接")
    add_title(ws, "重要链接汇总")
    links = [
        ["DeepSeek 控制台", "https://platform.deepseek.com", "注册 + 充值 + 获取Key"],
        ["Cherry Studio", "https://cherry-ai.com", "免费 AI 客户端下载"],
        ["OpenAI 控制台", "https://platform.openai.com", "Codex / ChatGPT API Key"],
        ["Anthropic 控制台", "https://console.anthropic.com", "Claude API Key"],
        ["Node.js 下载", "https://nodejs.org", "JavaScript 运行环境"],
        ["Claude Code 文档", "https://docs.anthropic.com/en/docs/claude-code", "Anthropic 官方文档"],
        ["Codex CLI GitHub", "https://github.com/openai/codex", "OpenAI 官方仓库"],
    ]
    style_sheet(ws, ["资源名称", "链接", "说明"], links, [28, 52, 40])

    xlsx_path = OUTPUT_DIR / "AI工具部署教程_小白版.xlsx"
    try:
        wb.save(str(xlsx_path))
    except PermissionError:
        xlsx_path = OUTPUT_DIR / "AI工具部署教程_小白版_v2.xlsx"
        wb.save(str(xlsx_path))
        print(f"  ⚠️ 原文件被占用，输出到: {xlsx_path.name}")
    print(f"✅ Excel: {xlsx_path} ({xlsx_path.stat().st_size//1024} KB, {len(wb.sheetnames)} sheets)")
    return xlsx_path


# ============================================================
# MAIN
# ============================================================
def main():
    print("=" * 55)
    print("  AI 部署教程 — PDF + PPT + Excel 生成器")
    print("=" * 55)

    print("\n📊 检查/生成图表...")
    charts = ensure_charts()
    for k, v in charts.items():
        if v and v.exists():
            print(f"  ✅ {v.name}")
        else:
            print(f"  ❌ {k}: MISSING")

    print("\n📄 生成 PDF...")
    pdf_path = build_pdf(charts)

    print("\n📊 生成 PPT...")
    from generate_ppt_v3 import build_ppt
    ppt_path = build_ppt(charts)

    print("\n📋 生成 Excel...")
    xlsx_path = build_excel()

    print("\n" + "=" * 55)
    print("  全部完成！")
    print(f"  📄 {pdf_path}")
    print(f"  📊 {ppt_path}")
    print(f"  📋 {xlsx_path}")
    print("=" * 55)

if __name__ == "__main__":
    main()
