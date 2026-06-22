#!/usr/bin/env python3
"""
AI工具部署教程 — PDF + PPT 生成器
面向小白，图文并茂，简单易懂
"""

import os
import sys
import textwrap
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent
IMG_DIR = OUTPUT_DIR / "images"
IMG_DIR.mkdir(exist_ok=True)

# ============================================================
# PART 1: 生成图表 (matplotlib)
# ============================================================

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Arc, Polygon
import numpy as np

# 中文字体设置
plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "WenQuanYi Micro Hei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

# 配色方案
C = {
    "primary": "#4F46E5",    # 靛蓝
    "success": "#059669",    # 绿
    "warning": "#D97706",    # 琥珀
    "danger": "#DC2626",     # 红
    "info": "#0891B2",       # 青
    "dark": "#1E293B",       # 深灰
    "light": "#F8FAFC",      # 浅灰
    "accent": "#7C3AED",     # 紫
    "bg": "#F1F5F9",         # 背景
}

def fig1_deployment_roadmap():
    """图1: 部署路线图 - 完整流程"""
    fig, ax = plt.subplots(1, 1, figsize=(10, 8))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")
    ax.set_facecolor("#FAFBFC")

    steps = [
        ("① 检查环境", "OS · 内存 · 网络", C["dark"], 1, 8.5),
        ("② 安装 Node.js", "nodejs.org下载LTS版", C["primary"], 3, 8.5),
        ("③ 打开终端", "Win: PowerShell\nMac: 终端.app", C["info"], 5, 8.5),
        ("④ 安装 Claude Code", "npm install -g\n@anthropic-ai/claude-code", C["accent"], 7, 8.5),
        ("⑤ 获取 API Key", "console.anthropic.com\n→ API Keys → 创建", C["warning"], 9, 8.5),
    ]

    steps2 = [
        ("⑥ 配置登录", "终端运行 claude login\n粘贴API Key", C["success"], 3, 4.5),
        ("⑦ 验证安装", "claude --version\nclaude \"你好\"", C["primary"], 6, 4.5),
        ("⑧ 开始使用！ 🎉", "可以开始和AI对话啦", C["danger"], 9, 4.5),
    ]

    all_steps = steps + steps2

    for title, desc, color, x, y in all_steps:
        # 圆角矩形盒子
        box = FancyBboxPatch((x-0.9, y-0.6), 1.8, 1.2,
                             boxstyle="round,pad=0.1",
                             facecolor=color, edgecolor="white",
                             linewidth=2, alpha=0.92)
        ax.add_patch(box)
        ax.text(x, y+0.2, title, ha="center", va="center",
                fontsize=10, fontweight="bold", color="white")
        ax.text(x, y-0.3, desc, ha="center", va="center",
                fontsize=7.5, color="white", alpha=0.9)

    # 连接箭头 (第一行)
    for i in range(4):
        ax.annotate("", xy=(all_steps[i+1][3]-0.95, all_steps[i+1][4]),
                    xytext=(all_steps[i][3]+0.95, all_steps[i][4]),
                    arrowprops=dict(arrowstyle="->", color=C["dark"], lw=2, connectionstyle="arc3,rad=0"))

    # 换行箭头 (步骤5→6)
    ax.annotate("", xy=(all_steps[5][3], all_steps[5][4]+0.7),
                xytext=(all_steps[4][3], all_steps[4][4]-0.7),
                arrowprops=dict(arrowstyle="->", color=C["dark"], lw=2, connectionstyle="arc3,rad=-0.3"))

    # 第二行箭头
    for i in range(5, 7):
        ax.annotate("", xy=(all_steps[i+1][3]-0.95, all_steps[i+1][4]),
                    xytext=(all_steps[i][3]+0.95, all_steps[i][4]),
                    arrowprops=dict(arrowstyle="->", color=C["dark"], lw=2, connectionstyle="arc3,rad=0"))

    ax.set_title("AI 工具部署完整路线图", fontsize=16, fontweight="bold", color=C["dark"], pad=15)

    path = IMG_DIR / "deployment_roadmap.png"
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white", edgecolor="none")
    plt.close(fig)
    print(f"  ✅ {path.name}")
    return path


def fig2_tool_comparison():
    """图2: 工具对比表 - 可视化柱状图"""
    fig, ax = plt.subplots(1, 1, figsize=(9, 5))
    ax.set_facecolor("#FAFBFC")

    tools = ["Claude Code\n(IDE插件)", "Claude Code\n(命令行)", "OpenClaw\n(龙虾)", "Stable\nDiffusion"]
    difficulties = [1, 2, 3, 3]
    times = [5, 15, 30, 60]
    colors_bar = [C["success"], C["primary"], C["warning"], C["danger"]]

    x = np.arange(len(tools))
    width = 0.35

    bars1 = ax.bar(x - width/2, difficulties, width, label="安装难度 (1-5)", color=colors_bar, alpha=0.85, edgecolor="white", linewidth=1.5)
    bars2 = ax.bar(x + width/2, [t/5 for t in times], width, label="预计时间 (÷5分钟)", color=colors_bar, alpha=0.3, edgecolor="white", linewidth=1.5, hatch="//")

    # 在柱子上标注实际分钟数
    for i, (bar, t) in enumerate(zip(bars2, times)):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.15,
                f"约{t}分钟", ha="center", fontsize=9, fontweight="bold", color=C["dark"])

    for i, (bar, d) in enumerate(zip(bars1, difficulties)):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height()/2,
                f"{'⭐'*d}", ha="center", fontsize=13)

    ax.set_ylabel("评级", fontsize=11, color=C["dark"])
    ax.set_xticks(x)
    ax.set_xticklabels(tools, fontsize=10)
    ax.set_ylim(0, 14)
    ax.legend(loc="upper left", fontsize=9, framealpha=0.9)
    ax.set_title("主流 AI 工具部署对比", fontsize=14, fontweight="bold", color=C["dark"], pad=12)
    ax.grid(axis="y", alpha=0.3, linestyle="--")

    # 推荐标签
    ax.annotate("👍 新手首选", xy=(0, difficulties[0]), xytext=(0.3, 4.5),
                fontsize=10, color=C["success"], fontweight="bold",
                arrowprops=dict(arrowstyle="->", color=C["success"], lw=2),
                bbox=dict(boxstyle="round,pad=0.3", facecolor="#D1FAE5", alpha=0.8))

    path = IMG_DIR / "tool_comparison.png"
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white", edgecolor="none")
    plt.close(fig)
    print(f"  ✅ {path.name}")
    return path


def fig3_troubleshooting():
    """图3: 常见问题排查流程图"""
    fig, ax = plt.subplots(1, 1, figsize=(10, 7.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")
    ax.set_facecolor("#FAFBFC")

    issues = [
        (5, 9.2, "🚫 npm: command not found", "Node.js 没装好", C["danger"],
         "→ 重新去 nodejs.org 下载 LTS 版本\n→ 安装时勾选「Add to PATH」\n→ 重新打开终端再试", 0),
        (5, 7.2, "🚫 安装报权限错误 (EACCES)", "权限不足", C["warning"],
         "Mac/Linux: 命令前加 sudo\nWindows: 用管理员身份运行终端", 1),
        (5, 5.2, "🚫 claude: command not found", "全局安装路径不在PATH", C["danger"],
         "→ 运行 npm list -g 检查是否安装成功\n→ 重启终端\n→ 如仍不行，用 npx @anthropic-ai/claude-code", 2),
        (5, 3.2, "🚫 API 返回 401 错误", "API Key 无效或过期", C["warning"],
         "→ 去 console.anthropic.com 检查 Key 状态\n→ 重新创建 Key 并配置\n→ 确认账户有余额", 3),
        (5, 1.2, "🚫 网络超时 / 连不上", "网络问题", C["info"],
         "→ 检查网络连接\n→ 关闭 VPN/代理试试\n→ 切换手机热点测试", 4),
    ]

    for x, y, title, cause, color, solution, idx in issues:
        box = FancyBboxPatch((0.5, y-0.7), 9, 1.4,
                             boxstyle="round,pad=0.15",
                             facecolor=color, edgecolor="white",
                             linewidth=2, alpha=0.12)
        ax.add_patch(box)

        # 左侧：标题 + 原因
        ax.text(1.2, y+0.15, title, fontsize=11, fontweight="bold", color=color, va="center")
        ax.text(1.2, y-0.35, f"原因: {cause}", fontsize=8.5, color="#64748B", va="center")

        # 右侧：解决方案
        ax.text(5.5, y+0.05, solution, fontsize=8.5, color=C["dark"], va="center",
                fontfamily="monospace", bbox=dict(boxstyle="round,pad=0.3", facecolor="white", alpha=0.8))

    ax.set_title("🔧 常见安装问题 & 解决方案速查", fontsize=15, fontweight="bold", color=C["dark"], pad=18)

    path = IMG_DIR / "troubleshooting.png"
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white", edgecolor="none")
    plt.close(fig)
    print(f"  ✅ {path.name}")
    return path


def fig4_env_check():
    """图4: 环境检查清单"""
    fig, ax = plt.subplots(1, 1, figsize=(9, 5))
    ax.set_facecolor("#FAFBFC")

    checks = [
        ("操作系统", "Win10+/Mac12+", "✅ 满足", C["success"]),
        ("内存", "≥ 8GB", "✅ 建议16GB", C["success"]),
        ("硬盘空间", "≥ 2GB可用", "✅ 足够", C["success"]),
        ("网络", "稳定宽带", "⚠️ 别用热点", C["warning"]),
        ("Node.js", "v18 LTS+", "📥 需安装", C["primary"]),
        ("终端/命令行", "PowerShell/Terminal", "✅ 系统自带", C["success"]),
        ("Anthropic账号", "需注册", "📝 需准备", C["accent"]),
        ("API Key", "需创建", "🔑 需准备", C["accent"]),
    ]

    y_positions = list(reversed(range(len(checks))))
    names = [c[0] for c in checks]
    requirements = [c[1] for c in checks]
    statuses = [c[2] for c in checks]
    colors_h = [c[3] for c in checks]

    bars = ax.barh(y_positions, [1]*len(checks), height=0.6, color=colors_h, alpha=0.2, edgecolor=colors_h, linewidth=1.5)

    for i, (name, req, status, color) in enumerate(zip(names, requirements, statuses, colors_h)):
        ax.text(0.02, y_positions[i], f"{name}", fontsize=11, fontweight="bold", color=C["dark"], va="center")
        ax.text(0.35, y_positions[i], f"要求: {req}", fontsize=9, color=C["dark"], va="center")
        ax.text(0.72, y_positions[i], status, fontsize=9, fontweight="bold", color=color, va="center",
                bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor=color, alpha=0.8))

    ax.set_yticks([])
    ax.set_xticks([])
    ax.set_xlim(0, 1.05)
    ax.set_title("📋 安装前环境检查清单", fontsize=14, fontweight="bold", color=C["dark"], pad=12)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_visible(False)
    ax.spines["left"].set_visible(False)

    path = IMG_DIR / "env_checklist.png"
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white", edgecolor="none")
    plt.close(fig)
    print(f"  ✅ {path.name}")
    return path


def fig5_install_steps_detail():
    """图5: 详细安装步骤图解 (4宫格)"""
    fig, axes = plt.subplots(2, 2, figsize=(11, 8))
    fig.suptitle("📖 Claude Code 详细安装四步法", fontsize=16, fontweight="bold", color=C["dark"], y=0.98)
    fig.set_facecolor("white")

    steps_data = [
        ("STEP 1: 安装 Node.js", [
            "① 浏览器打开 nodejs.org",
            "② 点击左侧 LTS 版本下载",
            "③ 双击 .msi 文件安装",
            "④ 一路「Next」到底",
            "⑤ ⚠️ 勾选 Add to PATH !",
            "⑥ 完成后打开终端验证:",
            "   $ node --version",
            "   v20.x.x ✅",
        ], C["success"]),
        ("STEP 2: 获取 API Key", [
            "① 浏览器打开:",
            "   console.anthropic.com",
            "② 注册/登录 Anthropic 账号",
            "③ 左侧菜单 → API Keys",
            "④ 点击 Create Key 按钮",
            "⑤ 复制 Key (只显示一次!)",
            "⑥ ⚠️ 妥善保存，不要泄露",
            "⑦ 账户需充值(最低$5)",
        ], C["warning"]),
        ("STEP 3: 安装 Claude Code", [
            "① 打开终端 (管理员模式)",
            "   Win: Win+X → 终端(管理员)",
            "   Mac: Cmd+Space → 终端",
            "② 运行安装命令:",
            "   $ npm install -g",
            "     @anthropic-ai/claude-code",
            "③ 等待下载完成(约2-5分钟)",
            "④ 验证: $ claude --version",
        ], C["primary"]),
        ("STEP 4: 配置 & 验证", [
            "① 终端运行: $ claude login",
            "② 粘贴你的 API Key",
            "③ 看到 ✓ 即配置成功",
            "④ 测试对话:",
            "   $ claude \"你好，介绍一下\"",
            "⑤ 看到回复 = 部署成功! 🎉",
            "⑥ 常用命令:",
            "   claude      进入对话",
            "   claude -p   单次提问",
        ], C["accent"]),
    ]

    for (title, lines, color), ax_i in zip(steps_data, axes.flat):
        ax_i.set_facecolor("#FAFBFC")
        ax_i.set_title(title, fontsize=12, fontweight="bold", color=color, pad=8)
        ax_i.axis("off")

        # 左边色条
        ax_i.add_patch(plt.Rectangle((0.02, 0.05), 0.03, 0.9, transform=ax_i.transAxes,
                                      facecolor=color, alpha=0.8, zorder=0))

        for j, line in enumerate(lines):
            is_code = line.startswith("   $") or line.startswith("   v")
            y_pos = 0.88 - j * 0.115
            if is_code:
                ax_i.text(0.1, y_pos, line, transform=ax_i.transAxes,
                         fontsize=7.5, color="#334155", fontfamily="monospace",
                         bbox=dict(boxstyle="round,pad=0.2", facecolor="#1E293B", alpha=0.08))
            elif line.startswith("①") or line.startswith("②") or line.startswith("③") or \
                 line.startswith("④") or line.startswith("⑤") or line.startswith("⑥") or line.startswith("⑦"):
                ax_i.text(0.08, y_pos, line, transform=ax_i.transAxes,
                         fontsize=8, color=C["dark"], fontweight="bold")
            elif line.startswith("⚠️"):
                ax_i.text(0.08, y_pos, line, transform=ax_i.transAxes,
                         fontsize=7.5, color=C["danger"])
            elif line.startswith("✓"):
                ax_i.text(0.08, y_pos, line, transform=ax_i.transAxes,
                         fontsize=8, color=C["success"], fontweight="bold")
            elif line.startswith("🎉"):
                ax_i.text(0.08, y_pos, line, transform=ax_i.transAxes,
                         fontsize=9, color=C["danger"], fontweight="bold")
            else:
                ax_i.text(0.08, y_pos, line, transform=ax_i.transAxes,
                         fontsize=8, color="#475569")

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    path = IMG_DIR / "install_steps.png"
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white", edgecolor="none")
    plt.close(fig)
    print(f"  ✅ {path.name}")
    return path


# ============================================================
# PART 2: 生成 PDF (weasyprint)
# ============================================================

def generate_pdf(img_paths):
    """从HTML生成精美PDF"""
    from weasyprint import HTML

    html_content = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<style>
  @page {{ margin: 2cm 1.8cm; size: A4; }}
  body {{
    font-family: "Microsoft YaHei", "SimHei", "PingFang SC", sans-serif;
    font-size: 13px; line-height: 1.8; color: #1E293B;
    background: #fff;
  }}
  .cover {{
    text-align: center; padding: 60px 0 40px 0;
    border-bottom: 4px solid #4F46E5; margin-bottom: 30px;
  }}
  .cover h1 {{
    font-size: 36px; color: #4F46E5; margin-bottom: 8px;
    letter-spacing: 2px;
  }}
  .cover h2 {{ font-size: 18px; color: #64748B; font-weight: normal; }}
  .cover .meta {{ color: #94A3B8; font-size: 12px; margin-top: 20px; }}

  h1 {{ font-size: 24px; color: #4F46E5; border-left: 5px solid #4F46E5; padding-left: 12px; margin-top: 40px; }}
  h2 {{ font-size: 18px; color: #1E293B; margin-top: 30px; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px; }}
  h3 {{ font-size: 15px; color: #334155; margin-top: 20px; }}

  .tip {{
    background: #EFF6FF; border-left: 4px solid #3B82F6;
    padding: 12px 18px; margin: 16px 0; border-radius: 4px;
    font-size: 12.5px;
  }}
  .warning {{
    background: #FFFBEB; border-left: 4px solid #F59E0B;
    padding: 12px 18px; margin: 16px 0; border-radius: 4px;
    font-size: 12.5px;
  }}
  .danger {{
    background: #FEF2F2; border-left: 4px solid #EF4444;
    padding: 12px 18px; margin: 16px 0; border-radius: 4px;
    font-size: 12.5px;
  }}
  .success {{
    background: #F0FDF4; border-left: 4px solid #22C55E;
    padding: 12px 18px; margin: 16px 0; border-radius: 4px;
    font-size: 12.5px;
  }}

  pre {{
    background: #1E293B; color: #E2E8F0; padding: 14px 18px;
    border-radius: 8px; font-size: 12px; overflow-x: auto;
    line-height: 1.6;
  }}
  code {{ background: #F1F5F9; padding: 2px 6px; border-radius: 3px; font-size: 12px; }}
  pre code {{ background: none; padding: 0; color: #E2E8F0; }}

  table {{ width: 100%; border-collapse: collapse; margin: 16px 0; font-size: 12px; }}
  th {{ background: #4F46E5; color: white; padding: 10px 14px; text-align: left; font-weight: 600; }}
  td {{ padding: 9px 14px; border-bottom: 1px solid #E2E8F0; }}
  tr:nth-child(even) {{ background: #F8FAFC; }}

  img.chart {{ width: 100%; max-width: 100%; margin: 20px 0; border-radius: 8px;
                box-shadow: 0 2px 12px rgba(0,0,0,0.08); }}

  .checklist {{ list-style: none; padding-left: 0; }}
  .checklist li::before {{ content: "☐ "; color: #4F46E5; font-weight: bold; }}
  .checklist li.done::before {{ content: "☑ "; color: #22C55E; }}

  .page-break {{ page-break-before: always; }}

  .toc {{ background: #F8FAFC; padding: 24px 30px; border-radius: 12px; margin: 20px 0; }}
  .toc a {{ color: #4F46E5; text-decoration: none; }}
  .toc ol {{ padding-left: 24px; }}
  .toc li {{ padding: 4px 0; }}
</style>
</head>
<body>

<!-- 封面 -->
<div class="cover">
  <h1>🤖 AI 工具部署教程</h1>
  <h2>从零开始，小白也能学会的 AI 工具安装指南</h2>
  <p class="meta">版本 1.0 · 2026-06-22 · 基于实战经验编写</p>
</div>

<!-- 目录 -->
<div class="toc">
  <h3 style="margin-top:0; color:#4F46E5;">📑 目录</h3>
  <ol>
    <li><a href="#ch1">认识 AI 工具部署</a> — 为什么需要，能做什么</li>
    <li><a href="#ch2">部署前的准备</a> — 环境检查，工具对比</li>
    <li><a href="#ch3">手把手安装教程</a> — 四步完成部署</li>
    <li><a href="#ch4">配置与验证</a> — 确保安装成功</li>
    <li><a href="#ch5">常见问题排查</a> — 报错不求人</li>
    <li><a href="#ch6">安装后使用指南</a> — 快速上手</li>
    <li><a href="#ch7">附录</a> — 速查表、资源链接</li>
  </ol>
</div>

<!-- 第一章 -->
<h1 id="ch1">一、认识 AI 工具部署</h1>

<h2>1.1 什么是 AI 工具部署？</h2>
<p><strong>AI 工具部署</strong> 就是把人工智能软件安装到你的电脑上，让它能正常运行。</p>
<p>类比：就像给手机安装微信、抖音 —— 只不过 AI 工具的安装稍微复杂一点，需要一些命令行操作。本教程就是帮你学会这一步。</p>

<h2>1.2 部署后能做什么？</h2>
<table>
  <tr><th>🎯 场景</th><th>💡 例子</th></tr>
  <tr><td>编程辅助</td><td>写代码、修Bug、解释代码逻辑</td></tr>
  <tr><td>文档处理</td><td>写报告、翻译、润色文章</td></tr>
  <tr><td>数据分析</td><td>处理Excel、生成图表、分析趋势</td></tr>
  <tr><td>学习助手</td><td>解答问题、整理笔记、生成练习题</td></tr>
  <tr><td>自动化工作流</td><td>自动发邮件、整理文件、定时任务</td></tr>
</table>

<div class="tip"><strong>💡 关键认知：</strong> AI 工具是你的「智能助手」，不是替代你，而是帮你省时间。安装好之后，你就拥有了一个 24 小时在线的 AI 帮手。</div>

<h2>1.3 部署路线图一览</h2>
<img class="chart" src="{img_paths['roadmap']}" alt="部署路线图">

<div class="page-break"></div>

<!-- 第二章 -->
<h1 id="ch2">二、部署前的准备</h1>

<h2>2.1 环境检查清单</h2>
<img class="chart" src="{img_paths['env']}" alt="环境检查">

<div class="warning"><strong>⚠️ 安装前必读：</strong> 请逐项核对上方清单。如果某项不满足，先解决再往下走。</div>

<h2>2.2 AI 工具选哪个？</h2>
<p>市面上有多种 AI 工具可选，以下是适合小白的对比：</p>

<img class="chart" src="{img_paths['compare']}" alt="工具对比">

<table>
  <tr><th>工具</th><th>合适人群</th><th>安装难度</th><th>推荐指数</th></tr>
  <tr><td>Claude Code (IDE插件)</td><td>会用 VSCode 的</td><td>⭐ 极简</td><td>⭐⭐⭐⭐⭐</td></tr>
  <tr><td>Claude Code (命令行)</td><td>不介意敲命令</td><td>⭐⭐ 简单</td><td>⭐⭐⭐⭐⭐</td></tr>
  <tr><td>OpenClaw (龙虾)</td><td>需要多功能集成</td><td>⭐⭐⭐ 中等</td><td>⭐⭐⭐⭐</td></tr>
  <tr><td>Stable Diffusion</td><td>需要 AI 绘画</td><td>⭐⭐⭐⭐ 较难</td><td>⭐⭐⭐</td></tr>
</table>

<div class="tip"><strong>👍 推荐：</strong> 本教程以 <strong>Claude Code 命令行版</strong> 为主线 —— 功能最强、适用最广、步骤最清晰。</div>

<div class="page-break"></div>

<!-- 第三章 -->
<h1 id="ch3">三、手把手安装教程</h1>

<h2>3.1 总览：四步完成部署</h2>
<img class="chart" src="{img_paths['steps']}" alt="安装步骤">

<h2>3.2 第一步：安装 Node.js 环境</h2>
<p>Claude Code 基于 Node.js 运行，所以需要先装它。</p>

<h3>macOS 用户</h3>
<pre><code># 方法一：官网下载 (推荐)
浏览器打开 https://nodejs.org → 下载 LTS 版本 → 双击安装

# 方法二：Homebrew (如果你装了)
brew install node@20</code></pre>

<h3>Windows 用户</h3>
<pre><code>1. 浏览器打开 https://nodejs.org
2. 点击左侧绿色按钮 "LTS" 下载 .msi 文件
3. 双击运行，一路点 Next
4. ⚠️ 关键一步：确保勾选 "Add to PATH"
5. 安装完成后，重启终端</code></pre>

<h3>验证安装</h3>
<pre><code># 打开终端，输入以下命令：
node --version
# 应该输出类似: v20.14.0

npm --version
# 应该输出类似: 10.7.0</code></pre>

<div class="success"><strong>✅ 看到版本号 = Node.js 安装成功！</strong> 继续下一步。</div>

<h2>3.3 第二步：获取 Anthropic API Key</h2>
<pre><code>1. 浏览器打开 https://console.anthropic.com
2. 注册账号 (用邮箱即可)
3. 登录后，左侧菜单点击 "API Keys"
4. 点击 "Create Key" 按钮
5. 复制生成的 Key (sk-ant-...开头)
6. ⚠️ 这个 Key 只显示一次，务必保存好！
7. 在 Billing 页面充值 (最低 $5，约 ¥36)</code></pre>

<div class="danger"><strong>🔒 安全警告：</strong> API Key 就像你的银行卡密码 —— 不要发给任何人，不要截图发朋友圈，不要提交到 GitHub！</div>

<div class="page-break"></div>

<h2>3.4 第三步：安装 Claude Code</h2>
<pre><code># 在终端中运行 (Windows 用管理员模式)：

npm install -g @anthropic-ai/claude-code

# 等待约 2-5 分钟，看到类似输出就成功了：
# + @anthropic-ai/claude-code@1.x.x
# added 150 packages in 45s</code></pre>

<div class="warning"><strong>⚠️ Windows 用户注意：</strong> 如果报权限错误 (EACCES/EPERM)，请用管理员身份运行终端：Win+X → 终端(管理员)。</div>

<h3>验证安装</h3>
<pre><code>claude --version
# 输出类似: 1.0.28  ← 看到版本号 = 安装成功！</code></pre>

<h2>3.5 第四步：配置登录</h2>
<pre><code># 方式一：交互式登录 (推荐)
claude login
# 按提示粘贴你的 API Key，看到 ✓ 即完成

# 方式二：环境变量方式
# Windows (PowerShell):
$env:ANTHROPIC_API_KEY = "sk-ant-你的key"

# macOS / Linux:
export ANTHROPIC_API_KEY="sk-ant-你的key"</code></pre>

<h3>🔑 最终验证</h3>
<pre><code>claude "你好，请用一句中文介绍你自己"

# 如果看到 AI 的回复 → 🎉 部署成功！</code></pre>

<div class="success"><strong>🎉 恭喜！你已经完成了 AI 工具的部署。</strong>现在你拥有一个可以通过命令行随时调用的 AI 助手。</div>

<div class="page-break"></div>

<!-- 第四章 -->
<h1 id="ch4">四、配置与验证</h1>

<h2>4.1 安装后检查清单</h2>
<ul class="checklist">
  <li class="done">Node.js 版本 ≥ 18</li>
  <li class="done">npm 全局包中有 @anthropic-ai/claude-code</li>
  <li class="done">claude --version 输出版本号</li>
  <li class="done">API Key 已配置 (claude login 成功)</li>
  <li class="done">能正常发起对话并收到回复</li>
  <li>了解基础命令用法</li>
  <li>知道如何更新和卸载</li>
</ul>

<h2>4.2 基础命令速查</h2>
<table>
  <tr><th>命令</th><th>作用</th><th>示例</th></tr>
  <tr><td><code>claude</code></td><td>进入交互对话模式</td><td>直接在终端输入</td></tr>
  <tr><td><code>claude "问题"</code></td><td>单次提问</td><td><code>claude "1+1=?"</code></td></tr>
  <tr><td><code>claude -p "问题"</code></td><td>纯文本输出</td><td><code>claude -p "写一首诗"</code></td></tr>
  <tr><td><code>claude --version</code></td><td>查看版本</td><td>检查是否需要更新</td></tr>
  <tr><td><code>claude login</code></td><td>重新登录/换Key</td><td>Key过期时使用</td></tr>
  <tr><td><code>npm update -g @anthropic-ai/claude-code</code></td><td>更新到最新版</td><td>建议每周检查</td></tr>
  <tr><td><code>npm uninstall -g @anthropic-ai/claude-code</code></td><td>卸载</td><td>不想用了就删</td></tr>
</table>

<h2>4.3 进阶：VSCode 插件版</h2>
<p>如果你使用 VSCode 编辑器，还有更简单的方式：</p>
<pre><code>1. 打开 VSCode
2. 左侧点击「扩展」(Ctrl+Shift+X)
3. 搜索 "Claude Code"
4. 点击「安装」
5. 安装后按提示登录 → 直接在编辑器里用 AI</code></pre>
<div class="tip"><strong>💡 插件版 vs 命令行版：</strong> 插件版集成在编辑器里，写代码更方便；命令行版可以在任何地方使用，更灵活。两个可以同时装。</div>

<div class="page-break"></div>

<!-- 第五章 -->
<h1 id="ch5">五、常见问题排查</h1>

<img class="chart" src="{img_paths['trouble']}" alt="问题排查">

<h2>5.1 详细解决方案</h2>

<h3>问题 1：npm: command not found</h3>
<div class="danger"><strong>原因：</strong> Node.js 没有正确安装，或未添加到系统 PATH。<br>
<strong>解决：</strong> 重新从 <code>nodejs.org</code> 下载 LTS 版本，安装时确保勾选 <strong>"Add to PATH"</strong>。安装后<strong>重新打开终端</strong>再试。</div>

<h3>问题 2：权限错误 (EACCES / EPERM)</h3>
<div class="warning"><strong>原因：</strong> 当前用户没有写入系统目录的权限。<br>
<strong>解决：</strong>
<ul>
  <li><strong>Mac/Linux：</strong> 命令前加 <code>sudo</code>，如 <code>sudo npm install -g @anthropic-ai/claude-code</code></li>
  <li><strong>Windows：</strong> 用管理员身份运行终端 (Win+X → 终端(管理员))</li>
</ul>
</div>

<h3>问题 3：claude: command not found</h3>
<div class="warning"><strong>原因：</strong> npm 全局包的 bin 目录不在 PATH 中。<br>
<strong>解决：</strong>
<ol>
  <li>运行 <code>npm list -g --depth=0</code> 确认包已安装</li>
  <li>重启终端</li>
  <li>如果还不行，用 <code>npx @anthropic-ai/claude-code</code> 代替 <code>claude</code></li>
</ol>
</div>

<h3>问题 4：API 返回 401 错误</h3>
<div class="danger"><strong>原因：</strong> API Key 无效、过期或账户欠费。<br>
<strong>解决：</strong>
<ol>
  <li>去 <code>console.anthropic.com</code> → API Keys 检查 Key 状态</li>
  <li>如果 Key 被删除/过期，创建新 Key 并重新 <code>claude login</code></li>
  <li>检查 Billing 页面确认账户有余额</li>
</ol>
</div>

<h3>问题 5：网络超时 / 连接失败</h3>
<div class="warning"><strong>原因：</strong> 网络问题或防火墙拦截。<br>
<strong>解决：</strong>
<ol>
  <li>检查网络是否正常（打开网页试试）</li>
  <li>关闭 VPN/代理软件再试</li>
  <li>切换手机热点测试（排除宽带问题）</li>
  <li>Windows：检查防火墙是否拦截了 Node.js</li>
</ol>
</div>

<div class="page-break"></div>

<!-- 第六章 -->
<h1 id="ch6">六、安装后使用指南</h1>

<h2>6.1 你的第一次 AI 对话</h2>
<pre><code># 简单问答
claude "今天天气怎么样？"

# 让它帮你做事
claude "帮我写一段 Python 代码，读取 CSV 文件并画折线图"

# 翻译
claude "把这段话翻译成英文：人工智能正在改变世界"

# 解释概念
claude "什么是机器学习？用简单的例子解释"</code></pre>

<h2>6.2 高效使用技巧</h2>
<table>
  <tr><th>技巧</th><th>说明</th><th>例子</th></tr>
  <tr><td>🎯 越具体越好</td><td>模糊的问题得到模糊的答案</td><td>❌ "写个方案" → ✅ "写一份针对奶茶店的抖音运营方案，预算5000元，目标增粉1000"</td></tr>
  <tr><td>📎 给上下文</td><td>告诉AI你的背景和目的</td><td>"我是新手，请用小学生能听懂的话解释"</td></tr>
  <tr><td>🔄 可以追问</td><td>不满意就让它改</td><td>"太长了，压缩到200字以内"</td></tr>
  <tr><td>📂 处理文件</td><td>Claude Code 可以读你的文件</td><td><code>claude "分析这个 Excel 的销售趋势"</code></td></tr>
  <tr><td>⏰ 定期更新</td><td>保持工具最新</td><td><code>npm update -g @anthropic-ai/claude-code</code></td></tr>
</table>

<h2>6.3 费用说明</h2>
<table>
  <tr><th>项目</th><th>费用</th><th>说明</th></tr>
  <tr><td>Claude Code 软件</td><td>🆓 免费</td><td>开源工具，永久免费</td></tr>
  <tr><td>Anthropic API</td><td>💰 按量付费</td><td>按输入/输出 token 计费，日常使用每月约 ¥30-200</td></tr>
  <tr><td>API 最低充值</td><td>$5 (约¥36)</td><td>够试用一段时间</td></tr>
</table>

<div class="tip"><strong>💰 省钱技巧：</strong> 日常使用选 <code>claude-haiku-4-5</code> 模型（便宜、够用），重任务时才用 <code>claude-sonnet-4-6</code>。在配置文件里设置默认模型即可。</div>

<h2>6.4 持续学习资源</h2>
<ul>
  <li>📖 <strong>官方文档：</strong> <code>docs.anthropic.com</code> — 最权威的参考</li>
  <li>📺 <strong>B站教程：</strong> 搜索 "Claude Code 教程" — 视频演示</li>
  <li>💬 <strong>社区讨论：</strong> Reddit r/ClaudeAI — 用户交流</li>
  <li>🔄 <strong>更新日志：</strong> <code>docs.anthropic.com/en/docs/claude-code/overview</code></li>
</ul>

<div class="page-break"></div>

<!-- 第七章 -->
<h1 id="ch7">七、附录</h1>

<h2>7.1 部署流程速查表</h2>
<table>
  <tr><th>步骤</th><th>做什么</th><th>关键命令</th><th>预计时间</th></tr>
  <tr><td>1</td><td>装 Node.js</td><td>去 nodejs.org 下载</td><td>5分钟</td></tr>
  <tr><td>2</td><td>获取 API Key</td><td>console.anthropic.com</td><td>5分钟</td></tr>
  <tr><td>3</td><td>安装 Claude Code</td><td><code>npm install -g @anthropic-ai/claude-code</code></td><td>3分钟</td></tr>
  <tr><td>4</td><td>配置登录</td><td><code>claude login</code></td><td>1分钟</td></tr>
  <tr><td>5</td><td>验证</td><td><code>claude "你好"</code></td><td>1分钟</td></tr>
  <tr><td colspan="3"><strong>总计</strong></td><td><strong>约15分钟</strong></td></tr>
</table>

<h2>7.2 重要链接</h2>
<table>
  <tr><th>资源</th><th>链接</th></tr>
  <tr><td>Node.js 下载</td><td>https://nodejs.org</td></tr>
  <tr><td>Anthropic 控制台</td><td>https://console.anthropic.com</td></tr>
  <tr><td>Claude Code 文档</td><td>https://docs.anthropic.com/en/docs/claude-code</td></tr>
  <tr><td>VSCode 下载</td><td>https://code.visualstudio.com</td></tr>
</table>

<h2>7.3 环境变量配置 (可选)</h2>
<p>如果你不想每次都用 <code>claude login</code>，可以设置系统环境变量：</p>

<h3>Windows</h3>
<pre><code># PowerShell (仅当前用户)
[Environment]::SetEnvironmentVariable(
    "ANTHROPIC_API_KEY",
    "sk-ant-你的key",
    "User"
)

# 或通过图形界面：
# 设置 → 系统 → 关于 → 高级系统设置 → 环境变量 → 新建</code></pre>

<h3>macOS / Linux</h3>
<pre><code># 添加到 ~/.zshrc 或 ~/.bashrc
echo 'export ANTHROPIC_API_KEY="sk-ant-你的key"' >> ~/.zshrc
source ~/.zshrc</code></pre>

<div class="warning"><strong>⚠️ 提醒：</strong> 不要在共享电脑上存 API Key！不要在公开的脚本里写 Key！</div>

<br>
<hr>
<p style="text-align:center; color:#94A3B8; font-size:12px;">
  📖 AI 工具部署教程 v1.0 · 2026-06-22 · 让每个人都能用上 AI
</p>

</body>
</html>"""

    pdf_path = OUTPUT_DIR / "AI工具部署教程_小白版.pdf"
    HTML(string=html_content).write_pdf(pdf_path)
    print(f"\n📄 PDF 已生成: {pdf_path}")
    return pdf_path


# ============================================================
# PART 3: 生成 PPT (python-pptx)
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_slide_bg(slide, color_hex):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(color_hex))

def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color="#1E293B", bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*hex_to_rgb(color))
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline_box(slide, left, top, width, height, lines, font_size=14, color="#334155"):
    """lines: list of (text, bold, color_override, size_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, clr, sz = line_data, False, color, font_size
        else:
            text = line_data[0]
            bold = line_data[1] if len(line_data) > 1 else False
            clr = line_data[2] if len(line_data) > 2 else color
            sz = line_data[3] if len(line_data) > 3 else font_size

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = RGBColor(*hex_to_rgb(clr))
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return tf

def add_card(slide, left, top, width, height, color_hex):
    """添加一个圆角矩形卡片"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color_hex))
    shape.line.fill.background()
    return shape

def generate_ppt(img_paths):
    """生成精美PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # === Slide 1: 封面 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_slide_bg(slide, "#4F46E5")

    add_text_box(slide, 1, 1.5, 11.3, 1.5, "🤖 AI 工具部署教程",
                 font_size=48, color="#FFFFFF", bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.2, 11.3, 1, "从零开始 · 小白也能学会的 AI 工具安装指南",
                 font_size=22, color="#C7D2FE", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.8, 11.3, 0.5, "Claude Code 命令行版 · 实战教程 · 2026版",
                 font_size=16, color="#A5B4FC", alignment=PP_ALIGN.CENTER)

    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.3), Inches(3.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*hex_to_rgb("#818CF8"))
    line.line.fill.background()

    # === Slide 2: 目录 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FFFFFF")

    add_text_box(slide, 0.8, 0.4, 5, 0.6, "📑 课程目录", font_size=32, color="#4F46E5", bold=True)

    toc_items = [
        ("01", "认识 AI 工具部署", "为什么需要 / 能做什么"),
        ("02", "部署前的准备", "环境检查 / 工具对比"),
        ("03", "手把手安装教程", "四步完成部署"),
        ("04", "配置与验证", "确保安装成功"),
        ("05", "常见问题排查", "报错不求人"),
        ("06", "安装后使用指南", "快速上手"),
        ("07", "附录", "速查表 / 资源链接"),
    ]

    for i, (num, title, desc) in enumerate(toc_items):
        y = 1.5 + i * 0.8
        # 编号圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*hex_to_rgb(C["primary"]))
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, 2.0, y-0.05, 5, 0.4, title, font_size=20, bold=True, color="#1E293B")
        add_text_box(slide, 2.0, y+0.35, 6, 0.3, desc, font_size=13, color="#64748B")

    # === Slide 3: 部署路线图 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FAFBFC")
    add_text_box(slide, 0.5, 0.2, 8, 0.6, "🗺️ 部署路线图", font_size=30, color="#4F46E5", bold=True)
    slide.shapes.add_picture(str(img_paths['roadmap']), Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))

    # === Slide 4: 环境检查 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FAFBFC")
    add_text_box(slide, 0.5, 0.2, 8, 0.6, "📋 安装前环境检查", font_size=30, color="#4F46E5", bold=True)
    slide.shapes.add_picture(str(img_paths['env']), Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.1))

    # === Slide 5: 工具对比 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FAFBFC")
    add_text_box(slide, 0.5, 0.2, 8, 0.6, "🔍 主流 AI 工具对比", font_size=30, color="#4F46E5", bold=True)
    slide.shapes.add_picture(str(img_paths['compare']), Inches(1.5), Inches(0.9), Inches(10.3), Inches(6.1))

    # === Slide 6-9: 四步安装法 ===
    step_slides = [
        ("STEP 1", "安装 Node.js 环境", C["success"], [
            "📥 浏览器打开 nodejs.org",
            "📥 下载 LTS 版本（左侧绿色按钮）",
            "📥 双击安装，一路 Next",
            "⚠️ 务必勾选「Add to PATH」！",
            "✅ 验证: 终端输入 node --version",
            "✅ 看到版本号 = 成功",
        ]),
        ("STEP 2", "获取 Anthropic API Key", C["warning"], [
            "🌐 浏览器打开 console.anthropic.com",
            "👤 注册 / 登录 Anthropic 账号",
            "🔑 左侧菜单 → API Keys → Create Key",
            "📋 复制 Key（sk-ant-... 开头）",
            "⚠️ Key 只显示一次，立即保存！",
            "💰 Billing 页面充值（最低 $5）",
        ]),
        ("STEP 3", "安装 Claude Code", C["primary"], [
            "💻 打开终端（管理员模式）",
            "⌨️ 运行: npm install -g @anthropic-ai/claude-code",
            "⏳ 等待 2-5 分钟下载",
            "✅ 验证: claude --version",
            "💡 如报权限错，用管理员终端",
        ]),
        ("STEP 4", "配置登录 & 验证", C["accent"], [
            "🔐 终端运行: claude login",
            "📋 粘贴你的 API Key",
            "✅ 看到 ✓ 即配置成功",
            "🧪 测试: claude \"你好\"",
            "🎉 收到回复 = 部署完成！",
        ]),
    ]

    for step_num, step_title, color, items in step_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_bg(slide, "#FFFFFF")

        # 左侧彩色面板
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
        panel.line.fill.background()

        add_text_box(slide, 0.8, 1.5, 3.2, 1, step_num, font_size=20, color="#FFFFFF", bold=True, alignment=PP_ALIGN.LEFT)
        add_text_box(slide, 0.8, 2.2, 3.2, 1.5, step_title, font_size=36, color="#FFFFFF", bold=True, alignment=PP_ALIGN.LEFT)

        # 右侧内容
        for j, item in enumerate(items):
            y = 1.0 + j * 0.9
            # 圆点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(y+0.12), Inches(0.18), Inches(0.18))
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
            dot.line.fill.background()

            is_warn = item.startswith("⚠️")
            is_success = item.startswith("✅") or item.startswith("🎉")
            clr = C["danger"] if is_warn else (C["success"] if is_success else C["dark"])
            add_text_box(slide, 5.6, y, 7, 0.5, item, font_size=18, color=clr, bold=is_warn or is_success)

    # === Slide 10: 安装步骤总览图 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FAFBFC")
    add_text_box(slide, 0.5, 0.15, 8, 0.6, "📖 四步法完整图解", font_size=30, color="#4F46E5", bold=True)
    slide.shapes.add_picture(str(img_paths['steps']), Inches(0.3), Inches(0.8), Inches(12.7), Inches(6.5))

    # === Slide 11: 常用命令速查 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FFFFFF")
    add_text_box(slide, 0.8, 0.3, 8, 0.6, "⌨️ 常用命令速查表", font_size=30, color="#4F46E5", bold=True)

    cmds = [
        ("claude", "进入交互对话模式"),
        ('claude "问题"', "单次提问"),
        ("claude --version", "查看版本号"),
        ("claude login", "重新登录 / 换Key"),
        ("npm update -g @anthropic-ai/claude-code", "更新到最新版"),
        ("npm uninstall -g @anthropic-ai/claude-code", "卸载"),
    ]

    for j, (cmd, desc) in enumerate(cmds):
        y = 1.3 + j * 0.85
        add_card(slide, 1, y, 5.5, 0.65, "#F1F5F9")
        add_text_box(slide, 1.2, y+0.1, 5.2, 0.45, cmd, font_size=15, color="#4F46E5", bold=True)
        add_text_box(slide, 7, y+0.1, 5.5, 0.45, desc, font_size=15, color="#475569")

    # 右侧提示
    add_card(slide, 9, 1.3, 3.5, 5, "#EFF6FF")
    add_text_box(slide, 9.3, 1.6, 3, 0.5, "💡 小技巧", font_size=20, bold=True, color="#4F46E5")
    add_multiline_box(slide, 9.3, 2.3, 3, 3.5, [
        "• 按 ↑ 键可以调出上一条命令",
        "• Tab 键自动补全命令",
        "• Ctrl+C 中断当前操作",
        "• 在 VS Code 里用更方便",
        "• 建议每周更新一次",
    ], font_size=13, color="#334155")

    # === Slide 12: 问题排查 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FAFBFC")
    add_text_box(slide, 0.5, 0.15, 8, 0.6, "🔧 常见问题排查指南", font_size=30, color="#4F46E5", bold=True)
    slide.shapes.add_picture(str(img_paths['trouble']), Inches(0.3), Inches(0.8), Inches(12.7), Inches(6.5))

    # === Slide 13: 使用技巧 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FFFFFF")
    add_text_box(slide, 0.8, 0.3, 10, 0.6, "💡 高效使用技巧", font_size=30, color="#4F46E5", bold=True)

    tips = [
        ("🎯", "越具体越好", "模糊问题 → 模糊答案", "❌ 「写个方案」\n✅ 「写奶茶店抖音运营方案，预算5000」"),
        ("📎", "给上下文", "告诉AI背景和目的", "「我是新手，请用小学生能懂的话解释」"),
        ("🔄", "可以追问", "不满意就让它改", "「太长了，压缩到200字」「换个风格」"),
        ("📂", "处理文件", "能读写你的本地文件", "「分析这个CSV的销售趋势」"),
        ("⏰", "定期更新", "保持工具最新", "npm update -g @anthropic-ai/claude-code"),
    ]

    for j, (icon, title, desc, example) in enumerate(tips):
        x = 0.5 + j * 2.5
        add_card(slide, x, 1.3, 2.2, 5.2, "#F8FAFC")
        add_text_box(slide, x+0.15, 1.5, 1.9, 0.5, f"{icon} {title}", font_size=17, bold=True, color=C["dark"])
        add_text_box(slide, x+0.15, 2.1, 1.9, 0.35, desc, font_size=11, color="#64748B")
        add_text_box(slide, x+0.15, 2.7, 1.9, 1.5, example, font_size=10, color="#334155")

    # === Slide 14: 费用说明 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#FFFFFF")
    add_text_box(slide, 0.8, 0.3, 8, 0.6, "💰 费用说明", font_size=30, color="#4F46E5", bold=True)

    fee_items = [
        ("Claude Code 软件", "🆓 完全免费", "开源工具，永久免费使用", C["success"]),
        ("Anthropic API", "💰 按量付费", "月均约 ¥30-200（日常使用）", C["warning"]),
        ("API 最低充值", "$5 (约¥36)", "够试用 1-2 周", C["primary"]),
    ]

    for j, (name, price, note, color) in enumerate(fee_items):
        y = 1.5 + j * 1.8
        add_card(slide, 1, y, 11, 1.4, "#F8FAFC")
        # 左边色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(0.08), Inches(1.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
        bar.line.fill.background()

        add_text_box(slide, 1.5, y+0.15, 4, 0.4, name, font_size=20, bold=True, color=C["dark"])
        add_text_box(slide, 6, y+0.15, 3, 0.4, price, font_size=22, bold=True, color=color)
        add_text_box(slide, 1.5, y+0.7, 8, 0.4, note, font_size=14, color="#64748B")

    # === Slide 15: 结束页 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, "#4F46E5")

    add_text_box(slide, 1, 1.8, 11.3, 1, "🎉 恭喜完成学习！", font_size=42, color="#FFFFFF", bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.0, 11.3, 1, "现在你已经掌握了 AI 工具的完整部署流程",
                 font_size=20, color="#C7D2FE", alignment=PP_ALIGN.CENTER)

    # 3个下一步
    next_steps = [
        ("💻", "打开终端", "运行你的第一条 AI 命令"),
        ("📖", "阅读官方文档", "docs.anthropic.com"),
        ("🔄", "持续更新", "保持工具最新版本"),
    ]
    for j, (icon, title, desc) in enumerate(next_steps):
        x = 2.5 + j * 3
        add_card(slide, x, 4.2, 2.5, 1.8, "rgba(255,255,255,0.15)")
        add_text_box(slide, x+0.2, 4.4, 2.1, 0.5, f"{icon} {title}", font_size=17, bold=True, color="#FFFFFF")
        add_text_box(slide, x+0.2, 5.0, 2.1, 0.5, desc, font_size=12, color="#C7D2FE")

    add_text_box(slide, 1, 6.5, 11.3, 0.5, "教程 v1.0 · 2026-06-22",
                 font_size=12, color="#818CF8", alignment=PP_ALIGN.CENTER)

    ppt_path = OUTPUT_DIR / "AI工具部署教程_小白版.pptx"
    prs.save(ppt_path)
    print(f"📊 PPT 已生成: {ppt_path}")
    return ppt_path


# ============================================================
# MAIN
# ============================================================

def main():
    print("=" * 60)
    print("  AI 部署教程 — PDF & PPT 生成器")
    print("=" * 60)

    print("\n📊 生成图表...")
    img_paths = {
        'roadmap': fig1_deployment_roadmap(),
        'compare': fig2_tool_comparison(),
        'trouble': fig3_troubleshooting(),
        'env': fig4_env_check(),
        'steps': fig5_install_steps_detail(),
    }

    print("\n📄 生成 PDF...")
    pdf_path = generate_pdf(img_paths)

    print("\n📊 生成 PPT...")
    ppt_path = generate_ppt(img_paths)

    print("\n" + "=" * 60)
    print("  ✅ 全部完成！")
    print(f"  📄 PDF: {pdf_path}")
    print(f"  📊 PPT: {ppt_path}")
    print("=" * 60)


if __name__ == "__main__":
    main()
