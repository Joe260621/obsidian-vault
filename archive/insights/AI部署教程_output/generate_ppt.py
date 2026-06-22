#!/usr/bin/env python3
"""PPT 生成器 - 使用 python-pptx"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path(__file__).parent
IMG_DIR = OUTPUT_DIR / "images"

# 配色
C = {
    "primary": "#4F46E5", "success": "#059669", "warning": "#D97706",
    "danger": "#DC2626", "info": "#0891B2", "dark": "#1E293B",
    "accent": "#7C3AED", "bg": "#F8FAFC", "light_primary": "#E0E7FF"
}

def hex_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*hex_rgb(color))

def tb(slide, l, t, w, h, text, size=18, color=C["dark"], bold=False, align=PP_ALIGN.LEFT):
    """添加文本框"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*hex_rgb(color))
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    return tf

def card(slide, l, t, w, h, color, alpha=1.0):
    """圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_rgb(color))
    shape.line.fill.background()
    return shape

def add_img(slide, path, l, t, w, h=None):
    if Path(path).exists():
        kw = {"width": Inches(w)}
        if h:
            kw["height"] = Inches(h)
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), **kw)


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    img = lambda n: IMG_DIR / n

    # === S1: 封面 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, C["primary"])
    tb(s, 1, 1.5, 11.3, 1.5, "AI 工具部署教程", 48, "#FFFFFF", True, PP_ALIGN.CENTER)
    tb(s, 1, 3.2, 11.3, 1, "从零开始 · 小白也能学会的 AI 工具安装指南", 22, "#C7D2FE", align=PP_ALIGN.CENTER)
    tb(s, 1, 4.6, 11.3, 0.5, "Claude Code 命令行版 · 实战教程 · 2026版", 16, "#A5B4FC", align=PP_ALIGN.CENTER)
    tb(s, 1, 5.8, 11.3, 0.5, "v1.0 · 2026-06-22", 12, "#818CF8", align=PP_ALIGN.CENTER)

    # === S2: 目录 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FFFFFF")
    tb(s, 0.8, 0.4, 8, 0.6, "课程目录", 32, C["primary"], True)
    items = [
        ("01", "认识 AI 工具部署", "为什么需要 / 能做什么"),
        ("02", "部署前的准备", "环境检查 / 工具对比"),
        ("03", "手把手安装教程", "四步完成部署"),
        ("04", "配置与验证", "确保安装成功"),
        ("05", "常见问题排查", "报错不求人"),
        ("06", "安装后使用指南", "快速上手"),
        ("07", "附录", "速查表 / 资源链接"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = 1.5 + i * 0.8
        c = card(s, 1.2, y, 0.55, 0.55, C["primary"])
        c.text_frame.paragraphs[0].text = num
        c.text_frame.paragraphs[0].font.size = Pt(16)
        c.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(s, 2.0, y-0.05, 6, 0.4, title, 20, C["dark"], True)
        tb(s, 2.0, y+0.35, 6, 0.3, desc, 13, "#64748B")

    # === S3: 部署路线图 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FAFBFC")
    tb(s, 0.5, 0.15, 8, 0.6, "部署路线图", 30, C["primary"], True)
    add_img(s, img("deployment_roadmap.png"), 0.3, 0.9, 12.7)

    # === S4: 环境检查 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FAFBFC")
    tb(s, 0.5, 0.15, 8, 0.6, "安装前环境检查", 30, C["primary"], True)
    add_img(s, img("env_checklist.png"), 0.5, 0.9, 12.3)

    # === S5: 工具对比 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FAFBFC")
    tb(s, 0.5, 0.15, 8, 0.6, "主流 AI 工具对比", 30, C["primary"], True)
    add_img(s, img("tool_comparison.png"), 2, 0.9, 9.3)

    # === S6-S9: 四步安装 ===
    steps = [
        ("STEP 1", "安装 Node.js 环境", C["success"],
         ["浏览器打开 nodejs.org", "下载 LTS 版本（左侧绿色按钮）",
          "双击安装，一路 Next", "务必勾选「Add to PATH」！",
          "验证: 终端输入 node --version", "看到版本号 = 成功"]),
        ("STEP 2", "获取 API Key", C["warning"],
         ["浏览器打开 console.anthropic.com", "注册/登录 Anthropic 账号",
          "左侧菜单 → API Keys → Create Key", "复制 Key (sk-ant-...开头)",
          "Key 只显示一次，立即保存！", "Billing 页面充值 (最低$5)"]),
        ("STEP 3", "安装 Claude Code", C["primary"],
         ["打开终端 (管理员模式)", "npm install -g @anthropic-ai/claude-code",
          "等待 2-5 分钟下载", "验证: claude --version", "如报权限错，用管理员终端"]),
        ("STEP 4", "配置 & 验证", C["accent"],
         ["终端运行: claude login", "粘贴你的 API Key", "看到 ✓ 即配置成功",
          '测试: claude "你好"', "收到回复 = 部署完成！"]),
    ]
    for step_num, step_title, color, items in steps:
        s = prs.slides.add_slide(blank)
        set_bg(s, "#FFFFFF")
        # 左侧色板
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hex_rgb(color)); bar.line.fill.background()
        tb(s, 0.8, 1.5, 3.2, 0.6, step_num, 20, "#FFFFFF", True)
        tb(s, 0.8, 2.2, 3.2, 1.2, step_title, 36, "#FFFFFF", True)
        for j, item in enumerate(items):
            y = 1.0 + j * 0.9
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(y+0.12), Inches(0.18), Inches(0.18))
            dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(*hex_rgb(color)); dot.line.fill.background()
            is_warn = "务必" in item or "只显示" in item or "权限" in item
            clr = C["danger"] if is_warn else C["dark"]
            tb(s, 5.6, y, 7, 0.5, item, 18, clr, is_warn)

    # === S10: 四步图解 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FAFBFC")
    tb(s, 0.5, 0.1, 8, 0.6, "四步法完整图解", 30, C["primary"], True)
    add_img(s, img("install_steps.png"), 0.3, 0.8, 12.7)

    # === S11: 命令速查 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 8, 0.6, "常用命令速查表", 30, C["primary"], True)
    cmds = [
        ("claude", "进入交互对话模式"),
        ('claude "问题"', "单次提问"),
        ("claude --version", "查看版本号"),
        ("claude login", "重新登录 / 换Key"),
        ("npm update -g @anthropic-ai/claude-code", "更新到最新版"),
        ("npm uninstall -g @anthropic-ai/claude-code", "卸载"),
    ]
    for j, (cmd, desc) in enumerate(cmds):
        y = 1.3 + j * 0.85
        card(s, 1, y, 5.5, 0.65, "#F1F5F9")
        tb(s, 1.2, y+0.1, 5.2, 0.45, cmd, 15, C["primary"], True)
        tb(s, 7, y+0.1, 5.5, 0.45, desc, 15, "#475569")
    card(s, 9, 1.3, 3.5, 5, "#EFF6FF")
    tb(s, 9.3, 1.5, 3, 0.5, "小技巧", 20, C["primary"], True)
    tips = ["按 ↑ 键调出上一条命令", "Tab 键自动补全", "Ctrl+C 中断操作",
            "VSCode 插件版更方便", "建议每周更新一次"]
    for j, tip in enumerate(tips):
        tb(s, 9.3, 2.2+j*0.5, 3, 0.4, f"• {tip}", 13, "#334155")

    # === S12: 问题排查 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FAFBFC")
    tb(s, 0.5, 0.1, 8, 0.6, "常见问题排查指南", 30, C["primary"], True)
    add_img(s, img("troubleshooting.png"), 0.3, 0.8, 12.7)

    # === S13: 使用技巧 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 10, 0.6, "高效使用技巧", 30, C["primary"], True)
    tip_cards = [
        ("越具体越好", "模糊问题→模糊答案", "❌「写个方案」\n✅「写奶茶店抖音运营方案」"),
        ("给上下文", "告诉AI背景和目的", "「我是新手，请用小学生\n能懂的话解释」"),
        ("可以追问", "不满意就让它改", "「太长了，压缩到200字」\n「换个风格」"),
        ("处理文件", "能读写你的本地文件", "「分析这个CSV的\n销售趋势」"),
        ("定期更新", "保持工具最新", "npm update -g\n@anthropic-ai/claude-code"),
    ]
    for j, (title, desc, example) in enumerate(tip_cards):
        x = 0.5 + j * 2.5
        card(s, x, 1.3, 2.2, 5.5, "#F8FAFC")
        tb(s, x+0.15, 1.5, 1.9, 0.5, title, 17, C["dark"], True)
        tb(s, x+0.15, 2.2, 1.9, 0.35, desc, 11, "#64748B")
        tb(s, x+0.15, 2.9, 1.9, 1.8, example, 10, "#334155")

    # === S14: 费用说明 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, "#FFFFFF")
    tb(s, 0.8, 0.3, 8, 0.6, "费用说明", 30, C["primary"], True)
    fees = [
        ("Claude Code 软件", "完全免费", "开源工具，永久免费使用", C["success"]),
        ("Anthropic API", "按量付费", "月均约 ¥30-200（日常使用）", C["warning"]),
        ("API 最低充值", "$5 (约¥36)", "够试用 1-2 周", C["primary"]),
    ]
    for j, (name, price, note, color) in enumerate(fees):
        y = 1.5 + j * 1.8
        card(s, 1, y, 11, 1.4, "#F8FAFC")
        bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(0.08), Inches(1.4))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = RGBColor(*hex_rgb(color)); bar2.line.fill.background()
        tb(s, 1.5, y+0.15, 4, 0.4, name, 20, C["dark"], True)
        tb(s, 6, y+0.15, 3, 0.4, price, 22, color, True)
        tb(s, 1.5, y+0.7, 8, 0.4, note, 14, "#64748B")

    # === S15: 结束页 ===
    s = prs.slides.add_slide(blank)
    set_bg(s, C["primary"])
    tb(s, 1, 1.8, 11.3, 1, "恭喜完成学习！", 42, "#FFFFFF", True, PP_ALIGN.CENTER)
    tb(s, 1, 3.0, 11.3, 1, "现在你已经掌握了 AI 工具的完整部署流程", 20, "#C7D2FE", align=PP_ALIGN.CENTER)
    nexts = [("打开终端", "运行第一条 AI 命令"), ("阅读文档", "docs.anthropic.com"), ("持续更新", "保持工具最新版")]
    for j, (t, d) in enumerate(nexts):
        x = 2.5 + j * 3
        card(s, x, 4.2, 2.5, 1.8, "#6366F1")
        tb(s, x+0.2, 4.4, 2.1, 0.5, t, 17, "#FFFFFF", True)
        tb(s, x+0.2, 5.1, 2.1, 0.5, d, 12, "#C7D2FE")
    tb(s, 1, 6.5, 11.3, 0.5, "教程 v1.0 · 2026-06-22", 12, "#818CF8", align=PP_ALIGN.CENTER)

    ppt_path = OUTPUT_DIR / "AI工具部署教程_小白版.pptx"
    prs.save(str(ppt_path))
    print(f"✅ PPT 已生成: {ppt_path}")
    return ppt_path


if __name__ == "__main__":
    build_ppt()
