#!/usr/bin/env python3
"""PPT生成器 v3 — 简洁版，无emoji，小字号，大文本框"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path(__file__).parent
IMG_DIR = OUTPUT_DIR / "images"

C = {
    "primary": "#4F46E5", "success": "#059669", "warning": "#D97706",
    "danger": "#DC2626", "info": "#0891B2", "dark": "#1E293B",
    "accent": "#7C3AED", "ds": "#1A56DB",
}

def hr(h): return tuple(int(h.lstrip("#")[i:i+2], 16) for i in (0, 2, 4))

def build_ppt(charts):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(s, c):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(*hr(c))

    def top_bar(s, color):
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hr(color)); bar.line.fill.background()

    # 多段落文本框
    def tbox(s, l, t, w, h, lines, colors=None, sizes=None, bolds=None):
        box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        for i, txt in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt
            sz = (sizes[i] if sizes and i < len(sizes) and sizes[i] else 12)
            p.font.size = Pt(sz)
            clr = (colors[i] if colors and i < len(colors) and colors[i] else C["dark"])
            p.font.color.rgb = RGBColor(*hr(clr))
            p.font.bold = (bolds[i] if bolds and i < len(bolds) and bolds[i] else False)
            p.font.name = "Microsoft YaHei"
        return tf

    def add_img(s, key, l, t, w, h=None):
        p = charts.get(key)
        if p and p.exists():
            kw = {"width": Inches(w)}
            if h: kw["height"] = Inches(h)
            s.shapes.add_picture(str(p), Inches(l), Inches(t), **kw)

    def left_panel(s, color, title_lines, title_sizes, subtitle_lines, sub_sizes):
        """左侧彩色面板 + 标题"""
        panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.06), Inches(4.5), Inches(7.44))
        panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(*hr(color)); panel.line.fill.background()
        all_lines = title_lines + [""] + subtitle_lines
        all_sizes = title_sizes + [6] + sub_sizes
        all_bolds = [True]*len(title_lines) + [False]*(1+len(subtitle_lines))
        tbox(s, 0.6, 1.0, 3.5, 3.5, all_lines, colors=["#FFFFFF"]*len(all_lines), sizes=all_sizes, bolds=all_bolds)

    def step_items(s, items, start_y=0.8, color=C["primary"], fs=11):
        """右侧步骤列表"""
        for j, item in enumerate(items):
            y = start_y + j * 0.72
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.1), Inches(y+0.06), Inches(0.1), Inches(0.1))
            dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(*hr(color)); dot.line.fill.background()
            tbox(s, 5.4, y-0.02, 7.5, 0.6, [item], sizes=[fs])

    # ================================================================
    # S1: 封面
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, C["primary"])
    tbox(s, 1, 1.8, 11.3, 1.4, ["AI 工具部署教程"], colors=["#FFFFFF"], sizes=[46], bolds=[True])
    tbox(s, 1, 3.4, 11.3, 0.6, ["从零开始，小白也能学会的 AI 安装指南"], colors=["#C7D2FE"], sizes=[20])
    tbox(s, 1, 4.2, 11.3, 0.5, ["覆盖 Claude Code / Codex / DeepSeek / Cherry Studio"], colors=["#A5B4FC"], sizes=[14])
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.1), Inches(3.3), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(*hr("#818CF8")); line.line.fill.background()

    # ================================================================
    # S2: 目录
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["primary"])
    tbox(s, 0.8, 0.4, 10, 0.6, ["目录"], colors=[C["primary"]], sizes=[28], bolds=[True])
    toc = ["一、认识 AI 工具部署", "二、选择你的部署路线", "三、部署前的准备",
           "四、DeepSeek + Cherry Studio（国内推荐）", "五、Claude Code 安装（海外）",
           "六、Codex CLI 安装", "七、用 DeepSeek API 驱动 Codex",
           "八、中转站方案：Claude Code 走 DeepSeek",
           "九、常见问题排查", "十、安装后使用指南", "附录：速查表 & 链接"]
    for i, t in enumerate(toc):
        tbox(s, 1.2, 1.3 + i * 0.5, 11, 0.4, [f"0{i+1}"[-2:] + "  " + t], sizes=[13])

    # ================================================================
    # S3: 路线对比图
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC"); top_bar(s, C["primary"])
    tbox(s, 0.5, 0.3, 10, 0.5, ["两条部署路线对比"], colors=[C["primary"]], sizes=[24], bolds=[True])
    add_img(s, "ds_vs_claude", 0.5, 1.0, 12.3)

    # ================================================================
    # S4: 环境检查 & 工具选择
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC"); top_bar(s, C["primary"])
    tbox(s, 0.5, 0.3, 10, 0.5, ["环境检查 & 工具选择"], colors=[C["primary"]], sizes=[24], bolds=[True])
    add_img(s, "env", 0.3, 1.0, 6.3)
    add_img(s, "compare", 6.8, 1.0, 6.2)

    # ================================================================
    # S5: DeepSeek 注册安装
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["ds"])
    left_panel(s, C["ds"], ["路线 B", "DeepSeek API", "+ Cherry Studio"], [24, 30, 24],
               ["国内直连, 无需外网", "手机号注册, 微信/支付宝充值", "约 10 分钟完成"], [12, 12, 12])
    step_items(s, [
        "1. platform.deepseek.com（国内直接打开）",
        "2. 手机号注册 -> 点 API Keys -> 创建新 Key -> 复制保存（只显示一次）",
        "3. 点充值 -> 微信/支付宝扫码 -> 最低 10 元",
        "4. cherry-ai.com 下载 Cherry Studio -> 安装",
        "5. 启动 Cherry Studio -> 设置 -> 模型服务 -> 添加 DeepSeek",
        "6. 填入刚才复制的 Key -> 保存 -> 回到对话框即可使用",
    ], 0.8, C["ds"], 11)

    # ================================================================
    # S6: Cherry Studio 特性
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["ds"])
    tbox(s, 0.8, 0.4, 10, 0.6, ["Cherry Studio：国产免费 AI 客户端"], colors=[C["ds"]], sizes=[24], bolds=[True])
    feats = [
        ("国内直连, 无需外网", "支持 DeepSeek / OpenAI / 智谱..."),
        ("文件对话", "拖入 PDF / Word / Excel 直接提问"),
        ("美观易用", "图形界面操作, 比命令行友好"),
        ("完全免费", "开源软件, 无任何收费"),
        ("多模型对比", "同时问多个模型, 比较答案"),
        ("提示词库 + Markdown", "内置提示词, 支持 Markdown 渲染"),
    ]
    for j, (t, d) in enumerate(feats):
        x = 0.6 + (j % 3) * 4.2; y = 1.3 + (j // 3) * 2.8
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.4))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(*hr("#F0F7FF")); card.line.fill.background()
        tbox(s, x+0.25, y+0.3, 3.3, 1.5, [t, d], colors=[C["ds"], "#475569"], sizes=[15, 11], bolds=[True, False])

    # ================================================================
    # S7: Claude Code 四步法图
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC"); top_bar(s, C["primary"])
    tbox(s, 0.5, 0.3, 10, 0.5, ["路线 A：Claude Code 四步安装法"], colors=[C["primary"]], sizes=[24], bolds=[True])
    add_img(s, "steps", 0.5, 1.0, 12.3)

    # ================================================================
    # S8: Claude Code 详细步骤
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["primary"])
    left_panel(s, C["primary"], ["路线 A", "Claude Code", "四步安装"], [18, 30, 22],
               ["基于 Node.js 的命令行工具", "需要外网 + Visa 支付", "约 15 分钟完成"], [11, 11, 11])
    step_items(s, [
        "1. nodejs.org 下载 LTS 版 -> 安装时勾选 Add to PATH -> 装完重启终端",
        "2. 终端输入 node --version，看到版本号说明装好了",
        "3. console.anthropic.com（需外网）注册 -> API Keys -> 创建 -> 复制 Key",
        "4. 终端输入 npm install -g @anthropic-ai/claude-code",
        "5. 终端输入 claude login -> 粘贴 Key",
        "6. 终端输入 claude \"你好\" -> 收到回复即成功",
    ], 0.8, C["primary"], 11)

    # ================================================================
    # S9: Codex CLI 安装
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, "#10A37F")
    left_panel(s, "#10A37F", ["路线 A-2", "Codex CLI", "OpenAI 出品"], [18, 30, 22],
               ["基于 Node.js 的命令行工具", "需外网 + Visa（同 Claude Code）", "可用 ChatGPT 同款 Key"], [11, 11, 11])
    step_items(s, [
        "1. nodejs.org 下载 LTS 版安装（已装的可跳过）",
        "2. platform.openai.com（需外网）-> API Keys -> 创建 -> 复制 Key",
        "3. npm install -g @openai/codex",
        "4. Win: $env:OPENAI_API_KEY = \"你的Key\"   Mac: export OPENAI_API_KEY=\"你的Key\"",
        "5. codex \"你好\" -> 收到回复即成功",
    ], 0.8, "#10A37F", 11)

    # ================================================================
    # S10: DeepSeek 驱动 Codex
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, "#059669")
    left_panel(s, "#059669", ["最强方案", "用 DeepSeek", "来跑 Codex"], [18, 30, 22],
               ["Codex 是车, DeepSeek 是油", "接口格式一样, 换个地址就行", "费用只要 OpenAI 的 1/10", "国内直连 + 微信充值"], [11, 11, 11, 11])
    step_items(s, [
        "1. platform.deepseek.com 注册 -> API Keys -> 创建 Key -> 复制（和前面一样）",
        "2. npm install -g @openai/codex（前面已装的可跳过）",
        "3. 打开终端，复制粘贴下面两行（把钥匙替换成你自己的）：",
        '   $env:OPENAI_BASE_URL = "https://api.deepseek.com/v1"',
        '   $env:OPENAI_API_KEY = "sk-粘贴你刚才复制的钥匙"',
        '4. 验证：codex "你好" -> 收到回复就是成功了！Codex 现在用 DeepSeek 驱动',
        "（Mac 用户把 $env: 改成 export，后面加等号）",
    ], 0.6, "#059669", 11)

    # ================================================================
    # S11: 中转站方案
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["accent"])
    tbox(s, 0.8, 0.4, 12, 0.5, ["进阶：中转站 -- 让 Claude Code 走 DeepSeek"], colors=[C["accent"]], sizes=[22], bolds=[True])
    tbox(s, 0.8, 1.0, 12, 0.5, ["Claude Code 不能直接连 DeepSeek（接口格式不同），需要一个中转站当翻译。整个过程在国内网络完成。"], sizes=[12])
    step_items(s, [
        "1. 找一个中转站（选支持微信/支付宝的、运营半年以上的、有客服群的）",
        "2. 注册后你会拿到两样东西：中转站地址（网址）和钥匙（sk-...）",
        "3. 打开终端，复制粘贴下面两行（替换成你自己的地址和钥匙）：",
        '   $env:ANTHROPIC_BASE_URL = "中转站给你的地址"',
        '   $env:ANTHROPIC_API_KEY = "中转站给你的钥匙"',
        '4. 正常使用：claude "你好" -> Claude Code 通过中转站走 DeepSeek 了',
        "选站提醒：先充 5-10 元试几天，稳定了再大额。别信白菜价宣传。",
    ], 1.8, C["accent"], 11)

    # ================================================================
    # S12: 问题排查
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FAFBFC"); top_bar(s, C["primary"])
    tbox(s, 0.5, 0.3, 10, 0.5, ["常见问题排查"], colors=[C["primary"]], sizes=[24], bolds=[True])
    add_img(s, "trouble", 0.3, 1.0, 12.7)

    # ================================================================
    # S13: 使用指南 + 命令速查
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["primary"])
    tbox(s, 0.8, 0.4, 10, 0.5, ["安装后使用指南"], colors=[C["primary"]], sizes=[24], bolds=[True])

    tips = [("越具体越好", "模糊问题 -> 模糊答案\n写方案 -> 写奶茶店运营方案"),
            ("给上下文", "告诉 AI 你的背景和目的\n我是新手, 用大白话解释"),
            ("可以追问", "不满意就让它改\n太长压缩到200字"),
            ("文件对话", "拖入 PDF/Excel 直接提问\n支持多种文件格式"),
            ("定期更新", "npm update -g 保持最新\nCherry Studio 自动更新")]
    for j, (t, d) in enumerate(tips):
        x = 0.5 + j * 2.5
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(2.2), Inches(3.0))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(*hr("#F8FAFC")); card.line.fill.background()
        tbox(s, x+0.15, 1.5, 1.9, 2.5, [t, "", d], colors=[C["dark"],"","#64748B"], sizes=[14, 4, 10], bolds=[True, False, False])

    tbox(s, 0.8, 4.8, 12, 1.8, [
        "常用命令：",
        "claude          进入 Claude Code 对话",
        "codex           进入 Codex 对话",
        'claude "问题"    单次提问（不进入对话模式）',
        'codex "问题"     单次提问',
        "claude/codex --version    查看版本号",
        "claude login    重新登录 / 更换 Key",
        "npm update -g @anthropic-ai/claude-code    更新 Claude Code",
        "npm update -g @openai/codex    更新 Codex",
    ], sizes=[13,11,11,11,11,11,11,11,11], colors=[C["dark"],"#475569","#475569","#475569","#475569","#475569","#475569","#475569","#475569"])

    # ================================================================
    # S14: 费用参考
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["primary"])
    tbox(s, 0.8, 0.4, 10, 0.5, ["费用参考"], colors=[C["primary"]], sizes=[24], bolds=[True])
    fees = [
        ("DeepSeek API", "1-2 元/百万 token", "极便宜, 国内直连, 微信充值", C["success"]),
        ("Codex + DeepSeek", "同上 (Codex 免费)", "Codex 免费 + DeepSeek 低价, 近乎免费", C["success"]),
        ("OpenAI API", "$0.15-15/百万 token", "需外网 + Visa, 小模型较便宜", C["info"]),
        ("Anthropic API", "$3-15/百万 token", "需外网 + Visa, 模型最强但贵", C["warning"]),
        ("Cherry Studio", "完全免费", "开源软件, 图形界面, 支持多模型", C["primary"]),
        ("Claude Code / Codex CLI", "完全免费", "开源命令行工具, npm 一键安装", C["primary"]),
    ]
    for j, (name, price, note, col) in enumerate(fees):
        y = 1.2 + j * 0.95
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y), Inches(0.05), Inches(0.75))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hr(col)); bar.line.fill.background()
        tbox(s, 1.3, y+0.02, 3, 0.3, [name], sizes=[14], colors=[C["dark"]], bolds=[True])
        tbox(s, 4.5, y+0.02, 4, 0.3, [price], sizes=[14], colors=[col], bolds=[True])
        tbox(s, 1.3, y+0.4, 10, 0.3, [note], sizes=[10], colors=["#64748B"])

    # ================================================================
    # S15: 重要链接
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, "#FFFFFF"); top_bar(s, C["primary"])
    tbox(s, 0.8, 0.4, 10, 0.5, ["重要链接"], colors=[C["primary"]], sizes=[24], bolds=[True])
    links = [
        ("DeepSeek 控制台", "platform.deepseek.com", "注册 + 获取 Key + 微信充值"),
        ("Cherry Studio", "cherry-ai.com", "下载免费 AI 客户端"),
        ("OpenAI 控制台", "platform.openai.com", "Codex / ChatGPT API Key"),
        ("Anthropic 控制台", "console.anthropic.com", "Claude API Key"),
        ("Node.js", "nodejs.org", "必须先装的运行环境"),
        ("Claude Code 文档", "docs.anthropic.com", "Claude Code 官方文档"),
        ("Codex CLI", "github.com/openai/codex", "OpenAI 官方开源仓库"),
    ]
    for j, (name, url, desc) in enumerate(links):
        y = 1.2 + j * 0.8
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(0.04), Inches(0.65))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hr(C["primary"])); bar.line.fill.background()
        tbox(s, 1.1, y+0.02, 3, 0.3, [name], sizes=[15], colors=[C["dark"]], bolds=[True])
        tbox(s, 4.2, y+0.02, 5, 0.3, [url], sizes=[12], colors=[C["primary"]])
        tbox(s, 1.1, y+0.38, 8, 0.25, [desc], sizes=[10], colors=["#64748B"])

    # ================================================================
    # S16: 结束页
    # ================================================================
    s = prs.slides.add_slide(blank); bg(s, C["primary"])
    tbox(s, 1, 1.8, 11.3, 1.2, ["恭喜完成学习！"], colors=["#FFFFFF"], sizes=[42], bolds=[True])
    tbox(s, 1, 3.2, 11.3, 0.8, ["现在你已经可以独立完成 AI 工具的部署"], colors=["#C7D2FE"], sizes=[18])
    nexts = [("小白首选", "DeepSeek + Cherry Studio"), ("编程首选", "Codex + DeepSeek"), ("全能王者", "Claude Code + 中转站")]
    for j, (t, d) in enumerate(nexts):
        x = 2 + j * 3.3
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.3), Inches(2.8), Inches(1.6))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(*hr("#6366F1")); card.line.fill.background()
        tbox(s, x+0.25, 4.5, 2.3, 0.8, [t, d], colors=["#FFFFFF", "#C7D2FE"], sizes=[16, 11], bolds=[True, False])

    ppt_path = OUTPUT_DIR / "AI工具部署教程_小白版.pptx"
    try:
        prs.save(str(ppt_path))
    except PermissionError:
        ppt_path = OUTPUT_DIR / "AI工具部署教程_小白版_v2.pptx"
        prs.save(str(ppt_path))
        print(f"  WARNING: locked, output to {ppt_path.name}")
    print(f"OK PPT: {ppt_path} ({ppt_path.stat().st_size//1024} KB, {len(prs.slides)} slides)")
    return ppt_path


if __name__ == "__main__":
    # standalone test
    charts = {k: IMG_DIR / f"{k}.png" for k in
        ["deployment_roadmap","tool_comparison","troubleshooting","env_checklist","install_steps","deepseek_vs_claude"]}
    build_ppt(charts)
